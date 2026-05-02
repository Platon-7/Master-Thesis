"""Compose RoboRef.pptx from the UvA template — fully native shapes/text.

Every visual element is an editable PowerPoint object: cards, color stripes,
bars, KPI tiles, score badges, trajectory dots, formulas. The only raster
content is the failsafe / DROID frame thumbnails on slides 6 and 7.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path("/gpfs/home3/pkarageorgis1/Master-Thesis/Robometer-LoRA/presentation")
TPL = ROOT / "uva-basis-powerpoint-template-1.pptx"
OUT = ROOT / "RoboRef.pptx"
ASSETS = ROOT / "assets"
RESULTS_PNG = Path("/gpfs/home3/pkarageorgis1/Master-Thesis/Robometer-LoRA/results/presentation")

# UvA palette
RED    = RGBColor(0xBC, 0x00, 0x31)
INK    = RGBColor(0x1F, 0x1D, 0x21)
GREY   = RGBColor(0x55, 0x55, 0x55)
LIGHT  = RGBColor(0xF4, 0xF1, 0xEE)   # SAND card background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x00, 0x4E, 0x92)
GREEN  = RGBColor(0x25, 0x78, 0x35)
ORANGE = RGBColor(0xE9, 0x83, 0x00)
PURPLE = RGBColor(0x75, 0x1B, 0x68)
YELLOW = RGBColor(0xBE, 0xB5, 0x11)
CYAN   = RGBColor(0x2A, 0xA5, 0xD0)
MUTE   = RGBColor(0x9C, 0xA3, 0xAF)

SCORE_PALETTE = {1: MUTE, 2: YELLOW, 3: ORANGE, 4: BLUE, 5: GREEN}

FONT = "Calibri"


# ---------------------------------------------------------------------------
# template helpers
# ---------------------------------------------------------------------------
def get_layout(prs, name):
    for L in prs.slide_layouts:
        if L.name == name:
            return L
    raise KeyError(name)


def clear_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def remove_picture_placeholders(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.type == 18:  # PICTURE
            sp = ph._element
            sp.getparent().remove(sp)


# ---------------------------------------------------------------------------
# text helpers
# ---------------------------------------------------------------------------
def set_runs(tf, runs, *, align=None, vertical_anchor=None, line_spacing=None):
    """Write a sequence of (text, kw) tuples into a text frame."""
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor
    p = tf.paragraphs[0]
    if align is not None: p.alignment = align
    if line_spacing is not None: p.line_spacing = line_spacing
    first = True
    for item in runs:
        if item == "\n":
            p = tf.add_paragraph()
            if align is not None: p.alignment = align
            if line_spacing is not None: p.line_spacing = line_spacing
            first = True
            continue
        text, kw = item
        if not first:
            pass
        run = p.add_run()
        run.text = text
        run.font.name = kw.get("font", FONT)
        run.font.size = Pt(kw.get("size", 12))
        run.font.bold = kw.get("bold", False)
        run.font.italic = kw.get("italic", False)
        run.font.color.rgb = kw.get("color", INK)
        first = False


def textbox(slide, x, y, w, h, runs, *, align=None,
            vertical_anchor=None, line_spacing=None, autofit=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    set_runs(tf, runs, align=align, vertical_anchor=vertical_anchor,
             line_spacing=line_spacing)
    return box


def simple_text(slide, x, y, w, h, text, **kw):
    return textbox(slide, x, y, w, h, [(text, kw)], align=kw.get("align"))


# ---------------------------------------------------------------------------
# shape helpers
# ---------------------------------------------------------------------------
def round_rect(slide, x, y, w, h, fill, *, line=None, line_w=0, radius_pct=0.04):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius_pct
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def rect(slide, x, y, w, h, fill, *, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    sh.shadow.inherit = False
    return sh


def oval(slide, cx, cy, d, fill, *, line=None, line_w=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(cx - d/2), Inches(cy - d/2),
                                Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def line(slide, x1, y1, x2, y2, color=INK, width=1.0):
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def card(slide, x, y, w, h, stripe_color, label, title, body,
         *, body_size=11, title_size=13, label_size=9, body_color=INK):
    """Sand-coloured rounded rectangle with a coloured left stripe."""
    round_rect(slide, x, y, w, h, LIGHT, radius_pct=0.05)
    rect(slide, x, y, 0.10, h, stripe_color)
    if label:
        simple_text(slide, x + 0.20, y + 0.10, w - 0.30, 0.30,
                    label.upper(), size=label_size, bold=True, color=GREY)
    simple_text(slide, x + 0.20, y + 0.40, w - 0.30, 0.55,
                title, size=title_size, bold=True, color=stripe_color)
    textbox(slide, x + 0.20, y + 1.00, w - 0.30, h - 1.10,
            [(body, dict(size=body_size, color=body_color))],
            line_spacing=1.15)


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def slide_title(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TITELDIA"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("RoboRef", dict(size=46, bold=True, color=INK)), "\n",
                ("A Foundation Reward Model", dict(size=28, color=INK)),
            ])
        elif idx == 16:
            set_runs(ph.text_frame, [
                ("MSc Artificial Intelligence  ·  Thesis defence",
                 dict(size=14, bold=True, color=RED)),
            ])
        elif idx == 1:
            set_runs(ph.text_frame, [
                ("Platon Karageorgis", dict(size=18, color=INK)),
            ])
        elif idx == 10:
            set_runs(ph.text_frame, [
                ("University of Amsterdam  ·  2026",
                 dict(size=12, color=GREY)),
            ])
    return slide


# ---------------------------------------------------------------------------
# Slide 2 — Outline
# ---------------------------------------------------------------------------
def slide_outline(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [("Outline", dict(size=32, bold=True, color=INK))])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    items = [
        ("01", "Robometer  —  the current state of the art"),
        ("02", "Dataset  —  what RBM-1M is, where it falls short"),
        ("03", "Contribution  —  dense per-frame failure annotation"),
        ("04", "Loss functions  —  three candidate objectives"),
        ("05", "Training strategy  —  LoRA bake-off"),
        ("06", "Results  —  test-set comparison and calibration"),
        ("07", "Future steps"),
    ]
    y = 2.20
    for n, label in items:
        textbox(slide, 0.80, y, 12.0, 0.55, [
            (f"{n}     ", dict(size=22, bold=True, color=RED)),
            (label, dict(size=20, color=INK)),
        ])
        y += 0.58
    return slide


# ---------------------------------------------------------------------------
# Slide 3 — Robometer overview (6 cards)
# ---------------------------------------------------------------------------
def slide_robometer(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Robometer  —  the current state of the art for general-purpose reward modelling",
                 dict(size=22, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "Aliang et al., 2026  ·  4-billion-parameter VLM reward model  ·  trained on RBM-1M (~1.7M trajectories).",
                size=12, color=GREY)

    cards = [
        ("Backbone", "Qwen3-VL-4B-Instruct", BLUE,
         "Multi-image vision–language transformer.  Inputs: trajectory frames + task instruction.  Outputs are pooled per frame into a per-frame hidden state."),
        ("Three task heads", "progress  ·  success  ·  preference", RED,
         "Progress: 10-bin C51 distribution over [0,1].  Success: per-frame binary logit.  Preference: pairwise ranking score for trajectory comparison."),
        ("Training signal", "Frame-level + trajectory-level", GREEN,
         "Progress on successes via t/T heuristic.  Failures supervised only by the preference head via paired comparisons.  No dense per-frame failure labels."),
        ("Dataset (RBM-1M)", "~1.7M trajectories  ·  93 archives", PURPLE,
         "Humanoid, human-hand, and standard arms across diverse embodiments and tasks.  Successes 8× more abundant than failures — heavy success bias."),
        ("Evaluation paradigm", "Reward alignment + policy ranking", ORANGE,
         "Per-frame Spearman / Kendall on labelled progress.  Ranking accuracy on held-out trajectory pairs across embodiments."),
        ("Open question for downstream RL", "False positives  →  reward hacking", CYAN,
         "Reward overestimation on failure trajectories breaks RL fine-tuning.  Robometer's preference-only failure supervision leaves this unaddressed."),
    ]
    cw, ch = 4.05, 2.55
    x0, y0 = 0.40, 2.20
    gap_x, gap_y = 0.10, 0.20
    for k, (lbl, title, col, body) in enumerate(cards):
        r, c = k // 3, k % 3
        x = x0 + c * (cw + gap_x)
        y = y0 + r * (ch + gap_y)
        card(slide, x, y, cw, ch, col, lbl, title, body,
             body_size=10.5, title_size=13, label_size=9)
    return slide


# ---------------------------------------------------------------------------
# Slide 4 — Dataset breakdown (native bar chart + KPI cards)
# ---------------------------------------------------------------------------
def slide_dataset(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Robometer's dataset  —  RBM-1M composition by embodiment family",
                 dict(size=22, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "1.69M trajectories from 93 archives  —  successes dominate; failures concentrate in standard arms.",
                size=12, color=GREY)

    # Bars
    families = [
        ("Humanoid",            551_147,      0,   8),
        ("Human / human-hand",  366_699,      0,  11),
        ("Standard robot arms", 558_476, 215_537, 74),
    ]
    bar_x = 2.50           # label width left of bar
    bar_y0 = 2.40
    bar_w_max = 6.40       # corresponds to total max episodes
    bar_h = 0.55
    row_h = 1.40
    max_total = max(s + f for _, s, f, _ in families)
    px_per_episode = bar_w_max / max_total

    for i, (name, succ, fail, archs) in enumerate(families):
        cy = bar_y0 + i * row_h
        # family label
        textbox(slide, 0.30, cy + (bar_h - 0.30) / 2, 2.10, 0.50,
                [(name, dict(size=13, color=INK))])

        # Success bar
        succ_w = succ * px_per_episode
        rect(slide, bar_x, cy, succ_w, bar_h, BLUE)
        if succ_w > 0.6:
            simple_text(slide, bar_x, cy + 0.12, succ_w, bar_h - 0.20,
                        f"{succ/1000:,.0f}k", size=12, bold=True,
                        color=WHITE, align=PP_ALIGN.CENTER)
        # Failure bar
        if fail > 0:
            fail_w = fail * px_per_episode
            rect(slide, bar_x + succ_w, cy, fail_w, bar_h, RED)
            if fail_w > 0.4:
                simple_text(slide, bar_x + succ_w, cy + 0.12, fail_w,
                            bar_h - 0.20, f"{fail/1000:,.0f}k",
                            size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # right-side annotations
        total = succ + fail
        right_x = bar_x + (total * px_per_episode) + 0.10
        textbox(slide, right_x, cy - 0.05, 1.90, 0.32,
                [(f"{total/1000:,.0f}k episodes", dict(size=11, bold=True, color=INK))])
        textbox(slide, right_x, cy + 0.27, 1.90, 0.30,
                [(f"{archs} archives", dict(size=10, color=GREY))])

    # X axis tick marks (visual only, no axis line)
    axis_y = bar_y0 + 3 * row_h - 0.60
    for tick_val in [0, 200_000, 400_000, 600_000, 800_000, 1_000_000]:
        tx = bar_x + tick_val * px_per_episode
        line(slide, tx, axis_y, tx, axis_y + 0.10, color=GREY, width=0.5)
        simple_text(slide, tx - 0.30, axis_y + 0.10, 0.60, 0.25,
                    f"{tick_val // 1000}k", size=9, color=GREY, align=PP_ALIGN.CENTER)
    line(slide, bar_x, axis_y, bar_x + bar_w_max + 0.05, axis_y, color=GREY, width=0.5)
    simple_text(slide, bar_x, axis_y + 0.40, bar_w_max, 0.30,
                "Episodes", size=10, color=GREY, align=PP_ALIGN.CENTER)

    # Legend
    legend_y = axis_y + 0.85
    rect(slide, bar_x, legend_y, 0.20, 0.20, BLUE)
    simple_text(slide, bar_x + 0.28, legend_y - 0.04, 1.80, 0.30,
                "Successful trajectories", size=11, color=INK)
    rect(slide, bar_x + 2.20, legend_y, 0.20, 0.20, RED)
    simple_text(slide, bar_x + 2.48, legend_y - 0.04, 1.80, 0.30,
                "Failure trajectories", size=11, color=INK)

    # KPI tiles on the right
    kpi_x = 9.40
    simple_text(slide, kpi_x, 2.10, 3.60, 0.30,
                "AT A GLANCE", size=10, bold=True, color=GREY)
    kpis = [
        ("1.69M",   "Total episodes"),
        ("1.48M",   "Successful trajectories"),
        ("215.5k",  "Failure trajectories"),
        ("93",      "Archives scanned"),
        ("68.9k",   "ICL pairs (failure → demo)"),
    ]
    ky = 2.45
    for val, lbl in kpis:
        round_rect(slide, kpi_x, ky, 3.60, 0.85, LIGHT, radius_pct=0.12)
        simple_text(slide, kpi_x + 0.18, ky + 0.13, 1.40, 0.60,
                    val, size=20, bold=True, color=RED)
        simple_text(slide, kpi_x + 1.58, ky + 0.27, 2.00, 0.40,
                    lbl, size=11, color=INK)
        ky += 0.95
    return slide


# ---------------------------------------------------------------------------
# Slide 5 — Contribution (two panels)
# ---------------------------------------------------------------------------
def slide_contribution(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Contribution  —  dense per-frame failure annotation at scale",
                 dict(size=24, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "Two complementary annotation channels populate failures with the same ordinal scale that successes already enjoy.",
                size=12, color=GREY)

    panels = [
        (BLUE,  "SIMULATOR-DERIVED   (Failsafe / MetaWorld)",
         "Procedural failure curriculum",
         [
             "Hand-crafted rubric injects 27 distinct failure modes per task.",
             "Three tasks (pick / push / stack) and three viewpoints (front / side / wrist).",
             "Per-frame label derived from simulator state (gripper pose, object-to-goal distance).",
             "Cleanest possible labels — used as the deciding evaluation split.",
         ],
         ["≈ 2,900 episodes  ·  3 tasks  ·  3 cameras",
          "Labels: {1, 2, 3, 4, 5}  —  no progress  →  success"]),
        (RED, "VLM + LLM   (DROID, Robometer Group A)",
         "Two-stage neural annotation pipeline",
         [
             "Stage 1 — Qwen3-VL describes each frame  (objects, gripper state, sub-step status).",
             "Stage 2 — Qwen3-LLM scores progress on the rubric using descriptions + task prompt.",
             "Decoupling vision from reasoning suppresses hallucinated rewards.",
             "Scales to ~10,500 real-world failure trajectories without human labelling.",
         ],
         ["≈ 5,500 DROID  +  ≈ 5,000 Robometer Group A",
          "Labels: {1, 2, 3, 4}  —  failures only, paired with success demos"]),
    ]
    pw = 6.05; ph_ = 4.85
    for i, (col, badge, headline, bullets, footer) in enumerate(panels):
        x = 0.50 + i * (pw + 0.20)
        y = 2.20
        round_rect(slide, x, y, pw, ph_, LIGHT, radius_pct=0.04)
        # Header band
        round_rect(slide, x, y, pw, 0.55, col, radius_pct=0.13)
        rect(slide, x, y + 0.30, pw, 0.25, col)  # square off bottom of band
        simple_text(slide, x, y + 0.10, pw, 0.40,
                    badge, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        simple_text(slide, x + 0.30, y + 0.75, pw - 0.5, 0.45,
                    headline, size=14, bold=True, color=col)
        # bullets
        by = y + 1.30
        for b in bullets:
            textbox(slide, x + 0.30, by, pw - 0.50, 0.65, [
                ("•   ", dict(size=12, color=col, bold=True)),
                (b, dict(size=11, color=INK)),
            ], line_spacing=1.15)
            by += 0.60
        # footer italic
        for j, f in enumerate(footer):
            simple_text(slide, x + 0.30, y + ph_ - 0.85 + j * 0.30,
                        pw - 0.5, 0.30, f, size=10.5, italic=True, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Slides 6 and 7 — frame strips + native trajectory line
# ---------------------------------------------------------------------------
def _trajectory(slide, xs_inch, y_top, y_bot, scores, score_max=5,
                line_color=RED, dot_d=0.18):
    """Native trajectory: connector lines + filled circles."""
    from pptx.enum.shapes import MSO_CONNECTOR
    n = len(scores)
    # value range
    vmin, vmax = 0.5, score_max + 0.5
    def vy(v):
        return y_bot - (v - vmin) / (vmax - vmin) * (y_bot - y_top)
    pts = list(zip(xs_inch, [vy(s) for s in scores]))
    # connectors
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln.line.color.rgb = line_color
        ln.line.width = Pt(2.0)
    # dots
    for (x, y), s in zip(pts, scores):
        oval(slide, x, y, dot_d, line_color, line=WHITE, line_w=1.0)
    # gridlines + y labels
    return vy


def slide_failsafe(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Failure annotation in simulation  —  Failsafe (ManiSkill, FailStackCube-v1)",
                 dict(size=22, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "Task:  Pick up the red cube and stack it on top of the green cube.",
                size=12, color=INK)
    simple_text(slide, 0.56, 2.10, 12.20, 0.30,
                "Failure scenario: grasp + carry away without dropping  ·  16 keyframes  ·  per-frame label from simulator state.",
                size=11, color=GREY)

    folder = ASSETS / "failsafe_example"
    files = sorted(folder.glob("frame_*.jpg"))[:16]
    scores = [1, 2, 3, 3, 3, 3, 3, 3, 3, 3, 3, 3, 3, 3, 3, 1]

    cols = 8
    img_w = 1.45; img_h = 1.00
    grid_x = 0.40; grid_y = 2.55
    gap = 0.07
    for k, (p, s) in enumerate(zip(files, scores)):
        r, c = k // cols, k % cols
        x = grid_x + c * (img_w + gap)
        y = grid_y + r * (img_h + 0.30 + gap)
        # frame title
        simple_text(slide, x, y, img_w, 0.20,
                    f"frame {k:02d}", size=9, color=GREY, align=PP_ALIGN.CENTER)
        # bordered image
        rect(slide, x - 0.025, y + 0.22, img_w + 0.05, img_h + 0.05,
             SCORE_PALETTE[s])
        slide.shapes.add_picture(str(p), Inches(x), Inches(y + 0.245),
                                 Inches(img_w), Inches(img_h))
        # score badge
        bx = x + img_w - 0.30; by = y + 0.27
        oval(slide, bx + 0.12, by + 0.12, 0.30, SCORE_PALETTE[s], line=WHITE, line_w=1.5)
        simple_text(slide, bx, by, 0.24, 0.24,
                    str(s), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # trajectory
    tj_x0, tj_x1 = 0.80, 12.50
    tj_y0, tj_y1 = 5.85, 6.85
    n = 16
    xs = [tj_x0 + (tj_x1 - tj_x0) * (i / (n - 1)) for i in range(n)]
    # axis ground line
    line(slide, tj_x0 - 0.05, tj_y1, tj_x1 + 0.05, tj_y1, color=GREY, width=0.5)
    # y labels
    labels = ["1  no progress", "2  approach", "3  grasp",
              "4  near completion", "5  success"]
    for i, lab in enumerate(labels, start=1):
        ly = tj_y1 - (i - 0.5) / 5.0 * (tj_y1 - tj_y0)
        simple_text(slide, 0.10, ly - 0.10, 0.65, 0.20,
                    str(i), size=10, color=GREY, align=PP_ALIGN.RIGHT)
    # gridlines
    for v in range(1, 6):
        ly = tj_y1 - (v - 0.5) / 5.0 * (tj_y1 - tj_y0)
        line(slide, tj_x0, ly, tj_x1, ly, color=RGBColor(0xE0, 0xE0, 0xE0), width=0.4)
    _trajectory(slide, xs, tj_y0, tj_y1, scores, score_max=5,
                line_color=RED, dot_d=0.16)
    # x labels
    for i in range(n):
        simple_text(slide, xs[i] - 0.20, tj_y1 + 0.05, 0.40, 0.22,
                    f"{i:02d}", size=8, color=GREY, align=PP_ALIGN.CENTER)
    simple_text(slide, tj_x0, tj_y1 + 0.30, tj_x1 - tj_x0, 0.25,
                "frame index", size=10, color=GREY, align=PP_ALIGN.CENTER)
    simple_text(slide, 0.80, 5.55, 8.0, 0.25,
                "Per-frame ground-truth ordinal label",
                size=11.5, bold=True, color=INK)
    return slide


def slide_vlm(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Failure annotation in the wild  —  DROID, two-stage VLM + LLM pipeline",
                 dict(size=22, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "Task:  Take cloth pieces out of container.", size=12, color=INK)
    simple_text(slide, 0.56, 2.10, 12.20, 0.30,
                "Pipeline:  Qwen3-VL describes each frame  →  Qwen3-LLM scores progress 1–4 on the rubric using description + task prompt.",
                size=11, color=GREY)

    folder = ASSETS / "droid_example"
    files = sorted(folder.glob("frame_*.jpg"))
    keep = [0, 2, 4, 6, 8, 10, 12, 14]
    scores = [1, 1, 2, 3, 3, 4, 4, 4]
    img_w = 1.50; img_h = 1.10
    grid_x = 0.40; grid_y = 2.55
    gap = 0.07
    for k, (i, s) in enumerate(zip(keep, scores)):
        x = grid_x + k * (img_w + gap)
        y = grid_y
        simple_text(slide, x, y, img_w, 0.20,
                    f"frame {i:02d}", size=9, color=GREY, align=PP_ALIGN.CENTER)
        rect(slide, x - 0.025, y + 0.22, img_w + 0.05, img_h + 0.05,
             SCORE_PALETTE[s])
        slide.shapes.add_picture(str(files[i]), Inches(x), Inches(y + 0.245),
                                 Inches(img_w), Inches(img_h))
        bx = x + img_w - 0.30; by = y + 0.27
        oval(slide, bx + 0.12, by + 0.12, 0.30, SCORE_PALETTE[s], line=WHITE, line_w=1.5)
        simple_text(slide, bx, by, 0.24, 0.24,
                    str(s), size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Trajectory
    tj_x0, tj_x1 = 0.80, 12.50
    tj_y0, tj_y1 = 4.30, 5.10
    n = len(scores)
    xs = [tj_x0 + (tj_x1 - tj_x0) * (i / (n - 1)) for i in range(n)]
    for v in range(1, 5):
        ly = tj_y1 - (v - 0.5) / 4.0 * (tj_y1 - tj_y0)
        line(slide, tj_x0, ly, tj_x1, ly, color=RGBColor(0xE0, 0xE0, 0xE0), width=0.4)
        simple_text(slide, 0.30, ly - 0.10, 0.45, 0.20,
                    str(v), size=10, color=GREY, align=PP_ALIGN.RIGHT)
    _trajectory(slide, xs, tj_y0, tj_y1, scores, score_max=4,
                line_color=RED, dot_d=0.18)
    for i, idx_v in enumerate(keep):
        simple_text(slide, xs[i] - 0.20, tj_y1 + 0.05, 0.40, 0.22,
                    f"{idx_v:02d}", size=9, color=GREY, align=PP_ALIGN.CENTER)
    simple_text(slide, 0.80, 4.00, 10.0, 0.25,
                "LLM-assigned ordinal progress score per keyframe",
                size=11.5, bold=True, color=INK)

    # VLM + LLM trace below
    round_rect(slide, 0.40, 5.65, 12.50, 1.55, LIGHT, radius_pct=0.04)
    rect(slide, 0.40, 5.65, 0.10, 1.55, RED)
    simple_text(slide, 0.60, 5.72, 12.20, 0.30,
                "Representative pipeline trace  ·  frame 10  ·  score 4",
                size=11.5, bold=True, color=INK)
    textbox(slide, 0.60, 6.00, 12.20, 0.55, [
        ("VLM description:  ", dict(size=10.5, bold=True, color=BLUE)),
        ("\"Wooden table with pink container holding cloth pieces. Robotic arm with gripper "
         "is positioned above the container. Cloth has been grasped and lifted, with one "
         "piece now resting on the table next to the container.\"",
         dict(size=10.5, color=INK))
    ], line_spacing=1.15)
    textbox(slide, 0.60, 6.70, 12.20, 0.45, [
        ("LLM verdict:  ", dict(size=10.5, bold=True, color=RED)),
        ("ANSWER: 4 because the robot has successfully removed the cloth from the container "
         "and placed it on the table, completing four of the six sub-steps.",
         dict(size=10.5, color=INK))
    ], line_spacing=1.15)
    return slide


# ---------------------------------------------------------------------------
# Slide 8 — Three loss functions
# ---------------------------------------------------------------------------
# Rich formula helper — uses DrawingML baseline shift for sub/super-scripts.
# This is what PowerPoint reliably renders inside text frames.
# ---------------------------------------------------------------------------
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _esc(text):
    return (text.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;"))


def formula_box(slide, x, y, w, h, lines, *, base_size=12, color=INK,
                align="ctr", line_spacing_pct=125):
    """Render a math-style textbox with sub/super-script-aware runs.

    lines: list of paragraphs.
    Each paragraph: list of (text, mode) tuples where mode ∈ {'', 'sub', 'sup', 'i'}.
        'i' = italic (variable letter convention).
    """
    from lxml import etree
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    txBody = tf._txBody
    for p in list(txBody.findall(f"{{{A_NS}}}p")):
        txBody.remove(p)

    color_hex = "%02X%02X%02X" % (color[0], color[1], color[2])
    sz_base = int(base_size * 100)
    sz_small = int(base_size * 75)

    for paragraph in lines:
        parts = [
            f'<a:p xmlns:a="{A_NS}">',
            f'  <a:pPr algn="{align}">',
            f'    <a:lnSpc><a:spcPct val="{line_spacing_pct * 1000}"/></a:lnSpc>',
            f'  </a:pPr>',
        ]
        for text, mode in paragraph:
            safe = _esc(text)
            italic = ' i="1"' if mode == "i" else ""
            if mode == "sub":
                rpr = (f'<a:rPr lang="en-US" sz="{sz_small}" baseline="-25000">'
                       f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
                       f'<a:latin typeface="{FONT}"/></a:rPr>')
            elif mode == "sup":
                rpr = (f'<a:rPr lang="en-US" sz="{sz_small}" baseline="30000">'
                       f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
                       f'<a:latin typeface="{FONT}"/></a:rPr>')
            else:
                rpr = (f'<a:rPr lang="en-US" sz="{sz_base}"{italic}>'
                       f'<a:solidFill><a:srgbClr val="{color_hex}"/></a:solidFill>'
                       f'<a:latin typeface="{FONT}"/></a:rPr>')
            parts.append(f'  <a:r>{rpr}<a:t>{safe}</a:t></a:r>')
        parts.append('</a:p>')
        txBody.append(etree.fromstring(''.join(parts)))
    return box


# ---------------------------------------------------------------------------
def slide_losses(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Three candidate objective functions",
                 dict(size=26, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.55, 12.20, 0.30,
                "Same backbone, same data, same LoRA setup  —  the only thing that varies is the loss applied to the heads.",
                size=12, color=GREY)

    titles = ["Robometer  (released baseline)",
              "Loss 1  ·  Asymmetric ordinal CORN",
              "Loss 2  ·  Asymmetric C51 + asymmetric BCE"]
    subtitles = [
        "Three independent heads  ·  failures supervised only via preference learning.",
        "Single ordinal head with cumulative thresholds  ·  over-prediction penalised harder.",
        "Both pretrained heads kept  ·  asymmetric weight damps over-confident success calls.",
    ]
    colors = [GREY, RED, BLUE]

    # Formulas as paragraphs of (text, mode) runs.  mode ∈ {'', 'sub', 'sup'}
    f_robometer = [
        [("L = L", ""), ("progress", "sub"),
         ("  +  L", ""), ("success", "sub"),
         ("  +  L", ""), ("pref", "sub")],
    ]
    f_loss1 = [
        # line 1
        [("L = −  Σ", ""), ("t,k", "sub"),
         ("    [   β", ""), ("k", "sub"),
         (" · b", ""), ("t,k", "sub"),
         (" · log σ(z", ""), ("t,k", "sub"), (")", "")],
        # line 2 (continuation, indented)
        [("           +  α", ""), ("k", "sub"),
         (" · (1 − b", ""), ("t,k", "sub"),
         (") · log(1 − σ(z", ""), ("t,k", "sub"), (")) ]", "")],
    ]
    f_loss2 = [
        [("L", ""), ("prog", "sub"),
         (" = w · CE(p, p", ""), ("*", "sup"),
         (")", "")],
        [("w = 𝟙[p̂ > p", ""), ("*", "sup"),
         ("]  +  λ · 𝟙[p̂ ≤ p", ""), ("*", "sup"), ("]", "")],
    ]
    formulas = [f_robometer, f_loss1, f_loss2]

    # Notation: list of (symbol_html_or_text, explanation)
    notations = [
        # Robometer
        [
            ("L^head", "loss applied to the corresponding head"),
            ("C51 CE", "cross-entropy over the 10-bin progress distribution"),
            ("BCE", "binary cross-entropy on the per-frame success logit"),
            ("rank", "pairwise margin loss (Robometer's preference learning)"),
        ],
        # Loss 1 — CORN
        [
            ("z_t,k", "k-th cumulative-threshold logit at frame t   (k = 2, 3, 4, 5)"),
            ("σ(z_t,k)", "= P(y_t ≥ k);   sigmoid of the threshold logit"),
            ("b_t,k", "1 if the true label y_t ≥ k, else 0"),
            ("β_k = 1", "uniform positive-class weight"),
            ("α_k = 1 + c·(k−2)", "negative weight grows with the threshold  →  false-success punished hardest"),
        ],
        # Loss 2
        [
            ("p̂", "model's expected progress for the frame  (continuous in [0, 1])"),
            ("p*", "target progress for the frame (rubric or t/T)"),
            ("CE(·, ·)", "Robometer's existing 10-bin discrete cross-entropy"),
            ("𝟙[·]", "indicator function"),
            ("λ ∈ (0, 1]", "damping factor for the under-prediction side  (default λ = 0.3)"),
        ],
    ]
    properties = [
        "Symmetric C51 / BCE  ·  no supervision on failure progress  ·  brittle as RL reward.",
        "Single ordinal head  ·  α_k > β_k makes it conservative  ·  P(success) = σ(z_5).",
        "Pretrained heads kept  ·  λ damps under-prediction  ·  no extra hyper-weights.",
    ]

    cw, ch = 4.05, 5.30
    x0, y0 = 0.40, 1.95
    gap = 0.10
    for i in range(3):
        x = x0 + i * (cw + gap); y = y0
        col = colors[i]
        round_rect(slide, x, y, cw, ch, LIGHT, radius_pct=0.04)
        rect(slide, x, y, 0.10, ch, col)

        simple_text(slide, x + 0.20, y + 0.15, cw - 0.30, 0.40,
                    titles[i], size=13, bold=True, color=col)
        textbox(slide, x + 0.20, y + 0.55, cw - 0.30, 0.55,
                [(subtitles[i], dict(size=10, color=INK))],
                line_spacing=1.15)

        # Head diagram band
        hd_y = y + 1.10
        if i == 0:
            heads = [("progress\nC51 × 10", BLUE),
                     ("success\nbinary", GREEN),
                     ("preference\npairwise", PURPLE)]
            bw = 1.10
            for j, (lbl, c2) in enumerate(heads):
                bx = x + 0.20 + j * (bw + 0.05)
                round_rect(slide, bx, hd_y, bw, 0.75, WHITE,
                           line=c2, line_w=1.5, radius_pct=0.10)
                simple_text(slide, bx, hd_y + 0.10, bw, 0.55,
                            lbl, size=9, bold=True, color=c2,
                            align=PP_ALIGN.CENTER)
        elif i == 1:
            round_rect(slide, x + 0.40, hd_y, cw - 0.80, 0.75, WHITE,
                       line=col, line_w=2.0, radius_pct=0.10)
            simple_text(slide, x + 0.40, hd_y + 0.10, cw - 0.80, 0.55,
                        "CORN ordinal head\n4 logits  ·  P(y ≥ k)",
                        size=10, bold=True, color=col, align=PP_ALIGN.CENTER)
        else:
            heads = [("progress\nC51 × 10\n(asym. CE)", BLUE),
                     ("success\nbinary\n(asym. BCE)", GREEN)]
            bw = 1.65
            for j, (lbl, c2) in enumerate(heads):
                bx = x + 0.30 + j * (bw + 0.10)
                round_rect(slide, bx, hd_y, bw, 0.75, WHITE,
                           line=c2, line_w=1.5, radius_pct=0.10)
                simple_text(slide, bx, hd_y + 0.05, bw, 0.65,
                            lbl, size=9, bold=True, color=c2,
                            align=PP_ALIGN.CENTER)

        # Math formula  —  one or two lines, slot is 0.85 inches tall
        formula_box(slide, x + 0.10, y + 1.95, cw - 0.20, 0.85,
                    formulas[i], base_size=11.5, color=col, align="ctr",
                    line_spacing_pct=115)

        # "where" notation block
        simple_text(slide, x + 0.20, y + 2.95, cw - 0.30, 0.25,
                    "where", size=9.5, bold=True, italic=True, color=col)
        ny = y + 3.18
        for sym, expl in notations[i]:
            textbox(slide, x + 0.20, ny, cw - 0.30, 0.40, [
                (sym + "  ", dict(size=10, bold=True, color=col)),
                ("— " + expl, dict(size=9.5, color=INK)),
            ], line_spacing=1.10)
            ny += 0.33

        # one-line property at the bottom of the card
        simple_text(slide, x + 0.20, y + ch - 0.40, cw - 0.30, 0.30,
                    properties[i], size=9.5, italic=True, color=GREY)
    return slide


# ---------------------------------------------------------------------------
# Slide 9 — Training strategy (6 cards)
# ---------------------------------------------------------------------------
def slide_training(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Training strategy", dict(size=26, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "LoRA fine-tune of Robometer-4B on ~18,900 balanced (failure | success) ICL pairs.",
                size=12, color=GREY)

    panels = [
        (RED,    "In-context learning  ·  per-example coin flip",
         "Each sample independently draws ICL on (prepend a success demo) or ICL off (query only) with p = 0.5. The demo defines what 'progress' means for this task."),
        (BLUE,   "Balanced batches  ·  50 / 50 success – failure",
         "Failure-query and success-query examples are constructed in equal numbers (~9.4k each). A stratified sampler enforces exact 50/50 within each batch after a short warmup."),
        (GREEN,  "LoRA adapters on Robometer-4B",
         "rank 32, α = 64, dropout 0.05.  Adapters on q/k/v/o + MLP gate/up/down. Backbone frozen, heads trained fully.  bf16 forward, fp32 adapters & heads."),
        (PURPLE, "Optimisation schedule",
         "AdamW · lr 1e-4 (adapters) · lr 5e-5 (heads) · weight decay 0.01.  Linear warmup over 5% of steps → cosine decay to 10% of peak.  7,500 steps · batch 8 · grad-clip 1.0 · seed 42."),
        (ORANGE, "Two-phase warmup",
         "First N steps draw failure-only batches to bootstrap the ordinal head before exposing the model to successes.  Loss 1: N = 2,000.  Loss 2: N = 1,000."),
        (CYAN,   "KL rehearsal anchor  ·  planned, full fine-tune",
         "FIFO buffer of past failure logits. On each success step sample one and add λ_KL · KL(P_old ∥ P_new) to prevent failure-prediction drift during success-heavy phases."),
    ]
    cw, ch = 4.05, 2.55
    x0, y0 = 0.40, 2.20
    gap_x, gap_y = 0.10, 0.20
    for k, (col, title, body) in enumerate(panels):
        r, c = k // 3, k % 3
        x = x0 + c * (cw + gap_x)
        y = y0 + r * (ch + gap_y)
        card(slide, x, y, cw, ch, col, "", title, body,
             body_size=10.5, title_size=12.5)
    return slide


# ---------------------------------------------------------------------------
# Slide 10 — Future steps
# ---------------------------------------------------------------------------
def slide_future(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Future steps", dict(size=32, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    items = [
        ("01", "Pick the winning loss",
         "Lock the LoRA bake-off after the held-out evaluation finishes; commit to one objective for the rest of the work."),
        ("02", "Explore further LoRA variants",
         "Sweep rank, ICL probability, label-smoothing on failures, and the asymmetric weights c and λ on the chosen loss."),
        ("03", "Pre-train rather than fine-tune",
         "Replace the LoRA adapter on Robometer-4B with a from-scratch run on the dense-label dataset to test whether the new supervision is enough on its own."),
        ("04", "GRPO with an unconstrained reward head",
         "Drop the rubric quantisation entirely and let GRPO shape an unbounded scalar reward; benchmark against the ordinal head."),
        ("05", "Submit to the RLC workshop",
         "Target the upcoming Reinforcement Learning Conference workshop track with the bake-off + ablations."),
        ("06", "Aim higher  ·  CoRL conference  (big bet)",
         "Lift the result from a workshop note to a full CoRL submission. Requires at least two real-world robot experiments end-to-end as the reward signal, on top of the bake-off and ablations."),
    ]
    y = 1.95
    for n, head, body in items:
        textbox(slide, 0.60, y, 12.20, 0.40, [
            (f"{n}    ", dict(size=16, bold=True, color=RED)),
            (head, dict(size=16, bold=True, color=INK)),
        ])
        textbox(slide, 1.40, y + 0.38, 11.40, 0.50,
                [(body, dict(size=11.5, color=INK))], line_spacing=1.15)
        y += 0.90
    return slide


# ---------------------------------------------------------------------------
# Native table helper
# ---------------------------------------------------------------------------
def native_table(slide, x, y, w, h, headers, rows, *,
                 header_fill=RED, header_color=WHITE, body_color=INK,
                 header_size=11, body_size=11.5, col_widths=None,
                 highlight_cells=None):
    """Construct a table with brand-aligned styling.

    highlight_cells: dict {(row_idx, col_idx): RGBColor} colours specific cell text.
    """
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols,
                                 Inches(x), Inches(y),
                                 Inches(w), Inches(h)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)
    # header
    for j, h_txt in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.margin_left = c.margin_right = Inches(0.10)
        c.margin_top = c.margin_bottom = Inches(0.06)
        set_runs(c.text_frame,
                 [(h_txt, dict(size=header_size, bold=True, color=header_color))],
                 align=PP_ALIGN.CENTER, vertical_anchor=MSO_ANCHOR.MIDDLE)
    # body
    highlight_cells = highlight_cells or {}
    for i, row in enumerate(rows, start=1):
        for j, cell_txt in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = LIGHT if i % 2 == 1 else WHITE
            c.margin_left = c.margin_right = Inches(0.10)
            c.margin_top = c.margin_bottom = Inches(0.06)
            color = highlight_cells.get((i - 1, j), body_color)
            bold = (i - 1, j) in highlight_cells
            align = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            set_runs(c.text_frame,
                     [(cell_txt, dict(size=body_size, bold=bold, color=color))],
                     align=align, vertical_anchor=MSO_ANCHOR.MIDDLE)
    return tbl


def add_image_centered(slide, png_path, x, y, w, h):
    """Insert an image fitted to (w, h) inches with aspect preserved, centered."""
    from PIL import Image
    iw, ih = Image.open(png_path).size
    ar_img = iw / ih
    ar_box = w / h
    if ar_img > ar_box:
        new_w = w; new_h = w / ar_img
    else:
        new_h = h; new_w = h * ar_img
    cx = x + (w - new_w) / 2
    cy = y + (h - new_h) / 2
    slide.shapes.add_picture(str(png_path), Inches(cx), Inches(cy),
                             Inches(new_w), Inches(new_h))


# ---------------------------------------------------------------------------
# Results 1 — Training dynamics
# ---------------------------------------------------------------------------
def slide_results_dynamics(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Training dynamics", dict(size=26, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.80, 12.20, 0.30,
                "Both runs converge cleanly  ·  10 evaluation rounds across 7,500 LoRA steps.",
                size=12, color=GREY)

    # Figure
    add_image_centered(slide,
        RESULTS_PNG / "fig_8_training_dynamics.png",
        0.40, 2.20, 12.50, 4.20)

    # Bottom takeaway band
    round_rect(slide, 0.40, 6.55, 12.50, 0.85, LIGHT, radius_pct=0.10)
    rect(slide, 0.40, 6.55, 0.10, 0.85, RED)
    textbox(slide, 0.65, 6.65, 12.20, 0.70, [
        ("L2 plateaus on eval ranking by step ~1,500  ·  ", dict(size=12.5, bold=True, color=BLUE)),
        ("the pre-trained C51 head finds its footing fast.\n",
         dict(size=12.5, color=INK)),
        ("L1's randomly-initialised CORN head climbs more irregularly through training.",
         dict(size=12.5, color=INK)),
    ], line_spacing=1.15)
    return slide


# ---------------------------------------------------------------------------
# Results 2 — Test-set comparison (3 metrics)
# ---------------------------------------------------------------------------
def slide_results_testset(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Test-set comparison  ·  Baseline ranks well but is unsafe; L2 is calibrated and deployable",
                 dict(size=20, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.85, 12.20, 0.30,
                "Each model evaluated on its trained success signal: success-head logit for Baseline and L2; σ(z_t,5) from CORN for L1.",
                size=12, color=GREY)

    # Table
    headers = ["Metric", "Baseline", "L1", "L2"]
    rows = [
        ["ROC-AUC (success)",                       "0.855  ✓", "0.651",      "0.783"],
        ["FPR @ τ = 0.5   (lower = safer)",          "0.195  ⚠",  "0.000",      "0.016  ✓"],
        ["ECE   (lower = better calibrated)",        "0.155",     "0.026",      "0.020  ✓"],
    ]
    highlight = {
        (0, 1): GREEN,   # baseline AUC winner
        (1, 3): GREEN,   # L2 FPR winner among useful models
        (1, 1): ORANGE,  # baseline FPR warning
        (2, 3): GREEN,   # L2 ECE winner
    }
    native_table(slide, 0.60, 2.30, 12.10, 1.85, headers, rows,
                 header_fill=RED, body_size=13, header_size=12,
                 col_widths=[4.30, 2.60, 2.60, 2.60],
                 highlight_cells=highlight)

    # Three takeaway bullets
    bullets = [
        (RED, "Baseline",
         "Best ranker but a 19.5% false-positive rate at the standard threshold — not deployable for RL "
         "without per-task threshold tuning."),
        (ORANGE, "L1",
         "Perfect FPR but ranks poorly (AUC 0.65) — its outputs are over-conservatised, "
         "almost never crossing 0.5  (revisited on the next slide)."),
        (GREEN, "L2",
         "The deployable winner: low FPR, calibrated probabilities, and competitive ranking."),
    ]
    by = 4.45
    for col, name, text in bullets:
        oval(slide, 0.75, by + 0.10, 0.20, col)
        textbox(slide, 1.00, by, 11.80, 0.60, [
            (f"{name}  ·  ", dict(size=13, bold=True, color=col)),
            (text, dict(size=12.5, color=INK)),
        ], line_spacing=1.15)
        by += 0.85

    return slide


# ---------------------------------------------------------------------------
# Results 3 — Recall at operating points
# ---------------------------------------------------------------------------
def slide_results_recall(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("How often does each model correctly identify a real success?",
                 dict(size=22, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.85, 12.20, 0.55,
                "A reward model says 'success' when its confidence exceeds a threshold τ.  Lower τ catches more "
                "successes but raises false alarms.  Below: TPR (% of real successes caught) at two practical "
                "operating points.",
                size=11.5, color=GREY)

    # Table
    headers = ["Operating point", "Baseline", "L1", "L2"]
    rows = [
        ["Safety-tuned   (only 5% false alarms allowed)",
         "catches 50% of successes",
         "catches 6%",
         "catches 62%  ✓"],
        ["Default threshold τ = 0.5   ('more likely than not')",
         "79% caught — but 20% false alarms  ⚠",
         "0% caught — model never confident enough to say 'success'",
         "61% caught, only 2% false alarms  ✓"],
    ]
    highlight = {
        (0, 3): GREEN,
        (1, 1): ORANGE,
        (1, 3): GREEN,
    }
    native_table(slide, 0.40, 2.55, 12.50, 1.40, headers, rows,
                 header_fill=RED, body_size=11, header_size=12,
                 col_widths=[3.80, 2.50, 3.40, 2.80],
                 highlight_cells=highlight)

    # Three model bullets
    bullets = [
        (GREY, "Baseline",
         "Ranks well but lies confidently. At τ = 0.5 it claims 'success' on 20% of real failures — "
         "that is reward hacking. Making it safe requires hand-tuning τ per task."),
        (ORANGE, "L1",
         "Over-corrected. The asymmetric loss made it so cautious that its outputs cap below 0.5 — even on "
         "real successes, L1 says 'I'm only 30–40% sure.'  Cannot separate successes from failures well."),
        (GREEN, "L2",
         "Hit the sweet spot. At τ = 0.5 it correctly identifies 61% of successes with only 2% false alarms; "
         "at the deployment-tuned threshold it catches 62%, beating Baseline's 50%."),
    ]
    by = 4.20
    for col, name, text in bullets:
        oval(slide, 0.75, by + 0.10, 0.20, col)
        textbox(slide, 1.00, by, 11.80, 0.65, [
            (f"{name}  ·  ", dict(size=12.5, bold=True, color=col)),
            (text, dict(size=11.5, color=INK)),
        ], line_spacing=1.15)
        by += 0.78

    # Bottom takeaway band
    round_rect(slide, 0.40, 6.65, 12.50, 0.65, LIGHT, radius_pct=0.12)
    rect(slide, 0.40, 6.65, 0.10, 0.65, GREEN)
    simple_text(slide, 0.65, 6.78, 12.20, 0.40,
                "L2 is the only model that is both safe and useful at standard deployment settings.",
                size=13, bold=True, color=INK)
    return slide


# ---------------------------------------------------------------------------
# Results 4 — Why L2 works (calibration)
# ---------------------------------------------------------------------------
def slide_results_calibration(prs):
    slide = prs.slides.add_slide(get_layout(prs, "TEKST"))
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_runs(ph.text_frame, [
                ("Why L2 works  ·  its predictions actually mean what they say",
                 dict(size=22, bold=True, color=INK))
            ])
        elif idx == 1:
            sp = ph._element; sp.getparent().remove(sp)

    simple_text(slide, 0.56, 1.85, 12.20, 0.55,
                "A model's confidence is meaningful only if it matches reality.  If the model says '70% confident' "
                "on 100 frames, ~70 of those should actually be successes.  ECE measures how far off the model is.",
                size=11.5, color=GREY)

    # Three big-number tiles
    tiles = [
        ("Baseline",  "0.155", "off by 15.5 percentage points on average", GREY),
        ("L1",        "0.026", "6× better than Baseline",                  ORANGE),
        ("L2",        "0.020", "8× better than Baseline  ✓",               GREEN),
    ]
    tile_w, tile_h = 4.05, 1.20
    tile_y = 2.50
    for i, (name, val, sub, col) in enumerate(tiles):
        x = 0.40 + i * (tile_w + 0.10)
        round_rect(slide, x, tile_y, tile_w, tile_h, LIGHT, radius_pct=0.06)
        rect(slide, x, tile_y, 0.10, tile_h, col)
        simple_text(slide, x + 0.20, tile_y + 0.10, tile_w - 0.30, 0.30,
                    name, size=11, bold=True, color=col)
        textbox(slide, x + 0.20, tile_y + 0.38, tile_w - 0.30, 0.55, [
            ("ECE = ", dict(size=14, color=INK)),
            (val, dict(size=24, bold=True, color=col)),
        ])
        simple_text(slide, x + 0.20, tile_y + 0.92, tile_w - 0.30, 0.25,
                    sub, size=10.5, italic=True, color=GREY)

    # Reliability diagram (image) + caption
    add_image_centered(slide,
        RESULTS_PNG / "fig_7_reliability_diagram.png",
        0.60, 3.85, 7.80, 2.55)
    simple_text(slide, 0.60, 6.45, 7.80, 0.25,
                "Each bar = a confidence bucket.  Diagonal = honest.  Below = overconfident.  Above = underconfident.",
                size=9.5, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    # Three model bullets on the right
    bullets = [
        (GREY,   "Baseline",
         "Lies confidently. At '90%' the actual success rate is closer to 50%. Outputs span [0,1] but are "
         "systematically overconfident."),
        (ORANGE, "L1",
         "Honest within a narrow range. Outputs cap at ~0.43, but inside that range what L1 says is what you "
         "get. Honesty bought at the cost of expressiveness."),
        (GREEN,  "L2",
         "Honest across its full range. Predictions hug the diagonal — predicted-30% means actually-30%. "
         "This is what makes it deployable."),
    ]
    bx = 8.65
    by = 3.95
    for col, name, text in bullets:
        oval(slide, bx + 0.10, by + 0.10, 0.18, col)
        textbox(slide, bx + 0.32, by, 4.20, 0.85, [
            (f"{name}  ·  ", dict(size=11.5, bold=True, color=col)),
            (text, dict(size=10.5, color=INK)),
        ], line_spacing=1.10)
        by += 0.85

    # Bottom takeaway band
    round_rect(slide, 0.40, 6.78, 12.50, 0.55, LIGHT, radius_pct=0.12)
    rect(slide, 0.40, 6.78, 0.10, 0.55, GREEN)
    simple_text(slide, 0.65, 6.86, 12.20, 0.40,
                "L2's outputs are the only ones you can plug into an RL system as continuous "
                "probabilities without re-calibrating per task.",
                size=12.5, bold=True, color=INK)
    return slide


# ---------------------------------------------------------------------------
def main():
    prs = Presentation(str(TPL))
    clear_all_slides(prs)

    slide_title(prs)
    slide_outline(prs)
    slide_robometer(prs)
    slide_dataset(prs)
    slide_contribution(prs)
    slide_failsafe(prs)
    slide_vlm(prs)
    slide_losses(prs)
    slide_training(prs)
    slide_results_dynamics(prs)
    slide_results_testset(prs)
    slide_results_recall(prs)
    slide_results_calibration(prs)
    slide_future(prs)

    prs.save(str(OUT))
    print(f"wrote {OUT}  ({OUT.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    main()
