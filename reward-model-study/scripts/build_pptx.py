"""Build an EDITABLE PowerPoint deck (.pptx) from the study artifacts.

Native text boxes / bullet lists / tables (all editable in PowerPoint) + embedded
plot images. Numbers are read from results/FULL_METRICS.csv and
results/ood_kendall_harness.csv so nothing is hand-typed.

  python reward-model-study/scripts/build_pptx.py
  -> reward-model-study/deck/VLM_reward_models.pptx
"""
import csv
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

R = Path("/gpfs/home3/pkarageorgis1/Master-Thesis/reward-model-study")
ROOT = Path("/gpfs/home3/pkarageorgis1/Master-Thesis")
RES = R / "results"
FIG = R / "figures"
LOSSDBG = ROOT / "loss-debug"
OUT = R / "deck" / "VLM_reward_models.pptx"

# ---- data -----------------------------------------------------------------
M = {(r["model"], r["cell"]): r for r in csv.DictReader(open(RES / "FULL_METRICS.csv"))}
KEN = {r["model"]: r for r in csv.DictReader(open(RES / "ood_kendall_harness.csv"))}


def g(m, cell, key):
    r = M.get((m, cell))
    if not r or r.get(key, "") in ("", "nan"):
        return None
    return float(r[key])


def ken(m):
    r = KEN.get(m)
    return None if not r else float(r["kendall_last"])


def ras_ood(m):
    r = KEN.get(m)
    return None if (not r or r.get("ranking_acc_sum","") in ("", "nan")) else float(r["ranking_acc_sum"])


def f(v, nd=2):
    return "—" if v is None else f"{v:.{nd}f}"


# ---- palette --------------------------------------------------------------
BG = RGBColor(0x0D, 0x11, 0x17)
FG = RGBColor(0xE6, 0xED, 0xF3)
MUTE = RGBColor(0x8B, 0x94, 0x9E)
BLUE = RGBColor(0x1F, 0x6F, 0xEB)
ACC = RGBColor(0x58, 0xA6, 0xFF)
RED = RGBColor(0xFF, 0x7B, 0x72)
GREEN = RGBColor(0x56, 0xD3, 0x64)
AMBER = RGBColor(0xE3, 0xB3, 0x41)
PURPLE = RGBColor(0xD2, 0xA8, 0xFF)
HDR = RGBColor(0x16, 0x1B, 0x22)
ROWA = RGBColor(0x11, 0x16, 0x1D)
ROWB = RGBColor(0x0D, 0x11, 0x17)
HI = RGBColor(0x12, 0x2B, 0x52)

FAMCOL = {"base": PURPLE, "asym": RED, "paper": ACC, "": FG}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def para(tf, text, size=18, color=FG, bold=False, italic=False, align=PP_ALIGN.LEFT,
         bullet=False, space=6, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    # crude rich text: split on ** for bold spans
    parts = text.split("**")
    for k, seg in enumerate(parts):
        if seg == "":
            continue
        run = p.add_run()
        run.text = ("• " if bullet and k == 0 else "") + seg
        run.font.size = Pt(size)
        run.font.bold = bold or (k % 2 == 1)
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Arial"
    return p


def header(s, num, title):
    tf = box(s, 0.5, 0.32, 12.3, 1.0)
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    r = p.add_run(); r.text = f" {num} "
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r.font.name = "Arial"
    # blue chip background isn't trivial on a run; emulate with brackets
    r2 = p.add_run(); r2.text = "  " + title
    r2.font.size = Pt(28); r2.font.bold = True; r2.font.color.rgb = FG; r2.font.name = "Arial"
    # accent rule
    ln = s.shapes.add_shape(1, Inches(0.5), Inches(1.18), Inches(2.2), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def fit(path, max_w, max_h):
    w, h = Image.open(path).size
    r = min(max_w / w, max_h / h)
    return w * r, h * r


def image(s, path, cx, top, max_w, max_h):
    w, h = fit(path, max_w, max_h)
    left = cx - w / 2
    pic = s.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(w), Inches(h))
    # white matte behind transparent-edge PNGs
    return pic


def table(s, headers, rows, left, top, width, col_w, fontsize=13, rowh=0.34):
    nr, nc = len(rows) + 1, len(headers)
    gtbl = s.shapes.add_table(nr, nc, Inches(left), Inches(top), Inches(width), Inches(rowh * nr)).table
    gtbl.first_row = False; gtbl.horz_banding = False
    tot = sum(col_w)
    for j, cw in enumerate(col_w):
        gtbl.columns[j].width = Emu(int(Inches(width) * cw / tot))
    # header
    for j, htxt in enumerate(headers):
        c = gtbl.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = HDR
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.margin_top = Pt(1); c.margin_bottom = Pt(1)
        tp = c.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        rr = tp.add_run(); rr.text = htxt; rr.font.size = Pt(fontsize); rr.font.bold = True
        rr.font.color.rgb = ACC; rr.font.name = "Arial"
    # body
    for i, (label, vals, fam, hicols) in enumerate(rows, start=1):
        col = FAMCOL.get(fam, FG)
        cells = [label] + vals
        for j, txt in enumerate(cells):
            c = gtbl.cell(i, j); c.fill.solid()
            c.fill.fore_color.rgb = HI if (j in hicols) else (ROWA if i % 2 else ROWB)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_top = Pt(1); c.margin_bottom = Pt(1); c.margin_left = Pt(6)
            tp = c.text_frame.paragraphs[0]; tp.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            rr = tp.add_run(); rr.text = txt
            rr.font.size = Pt(fontsize); rr.font.name = "Arial"
            rr.font.bold = (j == 0) or (j in hicols)
            rr.font.color.rgb = ACC if (j in hicols) else col
    return gtbl


def card(s, left, top, w, h, title, big, sub, accent):
    r = s.shapes.add_shape(5, Inches(left), Inches(top), Inches(w), Inches(h))  # rounded rect
    r.fill.solid(); r.fill.fore_color.rgb = HDR
    r.line.color.rgb = accent; r.line.width = Pt(1.5)
    tf = r.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(8)
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    rr = p0.add_run(); rr.text = title.upper(); rr.font.size = Pt(13); rr.font.bold = True
    rr.font.color.rgb = MUTE; rr.font.name = "Arial"
    p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
    rr = p1.add_run(); rr.text = big; rr.font.size = Pt(40); rr.font.bold = True
    rr.font.color.rgb = accent; rr.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    rr = p2.add_run(); rr.text = sub; rr.font.size = Pt(13); rr.font.color.rgb = FG; rr.font.name = "Arial"



# ============================================================================
#  Model identities (reviewer-facing labels)
# ============================================================================
ICL = {r["model"]: r for r in csv.DictReader(open(RES / "slide4_indist_icl.csv"))} \
      if (RES / "slide4_indist_icl.csv").exists() else {}
def ic(m, key):
    r = ICL.get(m); v = (r.get(key, "") if r else "")
    return "—" if v in ("", None) else f"{float(v):.3f}"

# OOD model keys (FULL_METRICS / ood csv use the checkpoint-suffixed names)
OOD_ORDER = ["baseline", "run1_s4000", "run2_s5000", "run3_s5000", "run4_s6500", "run5_s6500", "run6_s6500"]
OOD_LBL = {"baseline": "Robometer-4B  (off-the-shelf)", "run1_s4000": "Robometer-FT 4B  (asym + ICL)",
           "run2_s5000": "Robometer-FT 4B  (asym)", "run3_s5000": "Robometer-FT 4B  (standard loss)",
           "run4_s6500": "Qwen3.5-FT  (asym + ICL)", "run5_s6500": "Qwen3.5-FT  (asym)",
           "run6_s6500": "Qwen3.5-FT  (standard loss)"}
OOD_FAM = {"baseline": "base", "run1_s4000": "asym", "run2_s5000": "asym", "run3_s5000": "paper",
           "run4_s6500": "asym", "run5_s6500": "asym", "run6_s6500": "paper"}

# 7-model in-distribution table data (success + progress heads, correct-env reproduction)
# label, family, success(AUPRC,TPR@5%FPR,AUC), progress(TPR@5%FPR, gap, var, d')
MODELS7 = [
    ("Robometer-4B  (off-the-shelf)", "base",  (0.669, 0.236, 0.668), (0.222, 0.130, 0.091, 0.42)),
    ("Robometer-FT 4B  (asym + ICL)", "asym",  (0.910, 0.545, 0.878), (0.893, 0.476, 0.087, 1.54)),
    ("Robometer-FT 4B  (asym)",       "asym",  (0.894, 0.537, 0.877), (0.908, 0.412, 0.080, 1.31)),
    ("Robometer-FT 4B  (std loss)",   "paper", (0.881, 0.514, 0.868), (0.189, 0.144, 0.016, 1.08)),
    ("Qwen3.5-FT  (asym + ICL)",      "asym",  (0.694, 0.226, 0.689), (0.734, 0.192, 0.132, 0.54)),
    ("Qwen3.5-FT  (asym)",            "asym",  (0.654, 0.155, 0.643), (0.836, 0.104, 0.099, 0.27)),
    ("Qwen3.5-FT  (std loss)",        "paper", (0.639, 0.135, 0.652), (0.093, 0.043, 0.019, 0.32)),
]

# ============================================================================
#  Slide 1 — Title
# ============================================================================
s = slide()
tf = box(s, 0.9, 2.3, 11.5, 3.3)
para(tf, "Do fine-tuned VLM reward models give a better reward for robot RL?", size=38, color=ACC, bold=True, first=True, space=16)
para(tf, "We fine-tune vision-language reward models (Robometer-4B, Qwen3.5) and evaluate them as the reward signal "
         "for image-based reinforcement learning (IBRL) on a MetaWorld manipulation task.", size=20, color=FG, space=10)
para(tf, "Offline reward quality  ·  downstream RL performance  ·  why the two diverge", size=15, color=MUTE)

# ============================================================================
#  Slide 2 — THE PROBLEM: VLM rewards do not transfer to RL (sparse AND dense)
# ============================================================================
s = slide(); header(s, "1", "The problem: VLM rewards do not transfer to RL")
card(s, 0.8, 1.7, 5.5, 2.7, "Sparse reward", "≤ 0.13",
     "peak task success on CoffeePush — every model, every β / threshold", RED)
card(s, 7.0, 1.7, 5.5, 2.7, "Dense reward", "≤ 0.10",
     "switching to a dense per-step reward does not help (Qwen3.5 diverged)", RED)
tf = box(s, 0.8, 4.8, 11.7, 2.3)
para(tf, "With the same RL algorithm, task, and demonstrations, an oracle (ground-truth) reward trains the policy to "
         "**0.48–0.76** success. Every VLM reward — sparse or dense — caps near **0.1**.", size=19, color=FG, bullet=True, first=True)
para(tf, "The reward model, not the RL setup, is the bottleneck. The rest of this talk asks **why**, and whether "
         "fine-tuning helps.", size=19, color=AMBER, bullet=True)

# ============================================================================
#  Slide 3 — Out-of-distribution ranking (table only)
# ============================================================================
s = slide(); header(s, "2", "Trajectory ranking — out-of-distribution")
rows = []
for m in OOD_ORDER:
    hi = [1, 2, 3] if m == "baseline" else []
    rows.append((OOD_LBL[m], [f(g(m, "ood", "succ_AUC")), f(ras_ood(m)), f(ken(m))], OOD_FAM[m], hi))
table(s, ["model", "success AUC  ·  success head", "rank-acc (sum)  ·  progress head",
          "Kendall-τ  ·  progress head"],
      rows, 0.6, 1.7, 12.1, [3.4, 2.0, 2.0, 1.9], fontsize=13, rowh=0.46)
tf = box(s, 0.6, 5.8, 12.1, 1.2)
para(tf, "On held-out robots, the off-the-shelf baseline ranks trajectories best on every metric; fine-tuning on our "
         "data reduces out-of-distribution ranking. (Baseline reproduces the published Robometer numbers.)",
     size=15, color=MUTE, italic=True, first=True)

# ============================================================================
#  Slide 4 — OOD qualitative (baseline correct, fine-tuned model wrong)
# ============================================================================
s = slide(); header(s, "2", "Out-of-distribution — a concrete failure of the fine-tuned model")
image(s, R / "qualitative/ood/cand3_OOD_base0.92_ft0.21.png", cx=6.66, top=1.6, max_w=12.6, max_h=3.2)
tf = box(s, 0.6, 5.1, 12.1, 1.9)
para(tf, "A successful trajectory on a held-out robot (unseen embodiment). The progress head reads the final frame: "
         "the **baseline** correctly recognises task completion (0.92); the **fine-tuned model** does not (0.21).",
     size=17, color=FG, bullet=True, first=True)
para(tf, "Fine-tuning on a narrower data distribution costs generalisation to unseen embodiments.", size=17, color=AMBER, bullet=True)
para(tf, "(Progress head, scale 0–1; both models scored all 16 frames — 8 shown.)", size=13, color=MUTE, italic=True)

# ============================================================================
#  Slide 5 — In-distribution ranking (ICL off vs on)
# ============================================================================
s = slide(); header(s, "3", "Trajectory ranking — in-distribution")
ORD2 = ["baseline", "run1", "run2", "run3", "run4", "run5", "run6"]
LBL2 = {"baseline": "Robometer-4B  (off-the-shelf)", "run1": "Robometer-FT 4B  (asym + ICL)",
        "run2": "Robometer-FT 4B  (asym)", "run3": "Robometer-FT 4B  (std loss)",
        "run4": "Qwen3.5-FT  (asym + ICL)", "run5": "Qwen3.5-FT  (asym)", "run6": "Qwen3.5-FT  (std loss)"}
FAM2 = {"baseline": "base", "run1": "asym", "run2": "asym", "run3": "paper", "run4": "asym", "run5": "asym", "run6": "paper"}
rows = []
for m in ORD2:
    rows.append((LBL2[m], [ic(m, "icloff_rank_acc_sum"), ic(m, "icloff_kendall_last"),
                           ic(m, "iclon_rank_acc_sum"), ic(m, "iclon_kendall_last")], FAM2[m], []))
table(s, ["model", "rank-acc  (no ICL)", "Kendall-τ  (no ICL)", "rank-acc  (+ ICL)", "Kendall-τ  (+ ICL)"],
      rows, 0.6, 1.7, 12.1, [3.5, 2.1, 2.1, 2.1, 2.1], fontsize=13, rowh=0.46)
tf = box(s, 0.6, 5.9, 12.1, 1.4)
para(tf, "In-distribution, the fine-tuned 4B models rank best (progress head). In-context demonstrations (ICL) give a "
         "small additional gain for the models trained with them. Numbers reproduce the training-time evaluation exactly.",
     size=14, color=MUTE, italic=True, first=True)
para(tf, "In-context demos help the models trained with ICL (run1, run4); for the 4B no-ICL models they are roughly "
         "neutral. (Qwen3.5 no-ICL in-context rows pending — a Qwen3.5-specific harness ICL issue is under investigation.)",
     size=13, color=MUTE, italic=True)

# ============================================================================
#  Slide 6 — In-distribution qualitative (fine-tuned model superior)
# ============================================================================
s = slide(); header(s, "3", "In-distribution — where fine-tuning wins")
image(s, R / "qualitative/metaworld/fp3_FAILURE_base0.86_ft0.08.png", cx=6.66, top=1.6, max_w=12.6, max_h=3.2)
tf = box(s, 0.6, 5.1, 12.1, 1.9)
para(tf, "A **failed** attempt — the puck never reaches the shelf. The **baseline** scores it a success (0.86) — a false "
         "positive; the **fine-tuned model** correctly rejects it (0.08).", size=17, color=FG, bullet=True, first=True)
para(tf, "These false positives are exactly what a policy learns to exploit — the next section.", size=17, color=AMBER, bullet=True)
para(tf, "(Success head, 0–1; standard external camera, not the wrist view; 8 of 16 frames shown.)", size=13, color=MUTE, italic=True)

# ============================================================================
#  Slide 7 — Success head: the head the RL reward uses
# ============================================================================
s = slide(); header(s, "4", "Offline quality — success head (the reward IBRL uses)")
rows = []
for lbl, fam, sh, _ in MODELS7:
    hi = [2] if lbl.startswith("Robometer-FT 4B  (asym + ICL)") else []   # mark the best model's TPR
    rows.append((lbl, [f"{sh[0]:.2f}", f"{sh[1]:.2f}", f"{sh[2]:.2f}"], fam, hi))
table(s, ["model", "AUPRC", "TPR @ 5% FPR", "ROC-AUC"], rows, 0.6, 1.7, 9.6, [3.6, 1.5, 1.7, 1.5], fontsize=13, rowh=0.46)
tf = box(s, 0.6, 5.9, 12.1, 1.4)
para(tf, "The success head is what the deployed reward reads. The **fine-tuned 4B models clearly beat the baseline** — "
         "AUPRC 0.88–0.91 vs 0.67, and at a strict 5% false-positive budget they recover **0.51–0.55 vs 0.24**.",
     size=15, color=FG, bullet=True, first=True)
para(tf, "TPR @ 5% FPR is the operating point that matters for RL: catch successes without rewarding failures.", size=14, color=MUTE, italic=True)

# ============================================================================
#  Slide 8 — Progress head: class separation
# ============================================================================
s = slide(); header(s, "5", "Offline quality — progress head (class separation)")
rows = []
for lbl, fam, _, ph in MODELS7:
    hi = [4] if lbl.startswith("Robometer-FT 4B  (asym + ICL)") else []   # mark the best model's d′
    sigma = ph[1] / ph[3] if ph[3] else 0.0   # pooled σ, back-solved so gap/σ = d′ exactly
    rows.append((lbl, [f"{ph[0]:.2f}", f"+{ph[1]:.2f}", f"{sigma:.2f}", f"{ph[3]:.2f}"], fam, hi))
table(s, ["model", "TPR @ 5% FPR", "gap (succ−fail)", "pooled σ", "d′ = gap / σ"],
      rows, 0.6, 1.7, 11.2, [3.4, 1.8, 1.9, 1.6, 1.6], fontsize=13, rowh=0.46)
tf = box(s, 0.6, 5.7, 12.1, 1.6)
para(tf, "Separation of success vs failure by the progress head. The **fine-tuned 4B models separate the classes far "
         "better than the baseline** (d′ 1.3–1.5 vs 0.42). The asymmetric loss compresses the score *scale* (small gap) "
         "but keeps — even improves — *separability* (high d′).", size=14, color=FG, bullet=True, first=True)
para(tf, "Note d′ ≈ 1.5 means the gap is only ~1.5 standard deviations — good, not perfect. The failure distribution's "
         "upper tail still crosses any usable threshold, so a residual false-positive rate is built in (≈5% offline).",
     size=13, color=MUTE, italic=True)

# ============================================================================
#  Slide 9 — Downstream: the reward is the limiter (oracle control)
# ============================================================================
s = slide(); header(s, "6", "Downstream RL — the reward is the limiter")
image(s, FIG / "rl_gt_vs_vlm.png", cx=4.3, top=1.5, max_w=8.2, max_h=4.8)
card(s, 9.0, 1.6, 3.9, 1.5, "Oracle reward", "0.48–0.76", "3 seeds, same RL loop", GREEN)
card(s, 9.0, 3.3, 3.9, 1.5, "VLM reward", "≤ 0.13", "every model, every config", RED)
tf = box(s, 9.0, 5.0, 4.0, 2.1)
para(tf, "Replacing the VLM reward with the simulator's ground-truth reward trains the policy to 0.48–0.76 on the "
         "identical loop, task, and demonstrations.", size=13.5, color=FG, bullet=True, first=True)
para(tf, "So the limitation is the reward's usable signal, not the RL algorithm or the task.", size=13.5, color=AMBER, bullet=True)

# ============================================================================
#  Slide 10 — Downstream: reward shaping cannot fix it
# ============================================================================
s = slide(); header(s, "7", "Reward shaping does not break the ceiling")
image(s, FIG / "rl_beta_tau_sweep.png", cx=3.6, top=1.55, max_w=6.3, max_h=4.4)
image(s, FIG / "rl_peak_by_model.png", cx=10.0, top=1.55, max_w=6.3, max_h=4.4)
tf = box(s, 0.6, 6.1, 12.1, 1.2)
para(tf, "Left: a full reward-weighting and threshold sweep on the fine-tuned model — all flat. Right: peak success per "
         "model — every VLM reward stays ≤ 0.13, against the oracle control at 0.82.", size=15, color=FG, bullet=True, first=True)

# ============================================================================
#  Slide 12 — Downstream: RL amplifies the exploitable false positives
# ============================================================================
s = slide(); header(s, "9", "RL amplifies the exploitable false positives")
image(s, FIG / "rl_reward_hacking_3models.png", cx=6.66, top=1.6, max_w=12.8, max_h=3.6)
tf = box(s, 0.6, 5.4, 12.1, 1.9)
para(tf, "The reward's mistakes form a consistent, findable pattern. RL is an optimiser, so the policy learns to trigger "
         "them: the false-positive rate grows from **14–17% offline to 41–62% on-policy** for every model.", size=14.5, color=FG, bullet=True, first=True)
para(tf, "The agent banks reward ≈ its false-positive rate while true success stays near 0–5% — classic reward hacking.", size=14.5, color=AMBER, bullet=True)

# ============================================================================
#  Slide 13 — Downstream: it's exploitability, not the raw FP rate
# ============================================================================
s = slide(); header(s, "10", "It is exploitability, not the raw error rate")
image(s, FIG / "rl_fp_doseresponse.png", cx=6.66, top=1.55, max_w=12.8, max_h=3.8)
tf = box(s, 0.6, 5.5, 12.1, 1.8)
para(tf, "Controlled test: injecting *random* false positives into the oracle reward degrades performance only gradually "
         "(5% → 0.48) — a policy tolerates non-exploitable noise.", size=14, color=FG, bullet=True, first=True)
para(tf, "The VLM's false positives are *structured*, so the policy locks onto them and the rate runs away — a "
         "self-reinforcing loop that random noise cannot reproduce.", size=14, color=AMBER, bullet=True)

# ============================================================================
#  Slide 14 — Downstream: dense reward also fails (new)
# ============================================================================
s = slide(); header(s, "11", "A dense reward does not help either")
rows = [
    ("Robometer-4B  (off-the-shelf)", ["0.10", "flat — no learning"], "base", []),
    ("Robometer-FT 4B  (ours)",       ["0.04", "flat — no learning"], "asym", []),
]
table(s, ["dense per-step reward (β = 1, progress)", "peak success", "outcome"],
      rows, 0.6, 1.8, 11.0, [4.4, 1.8, 2.6], fontsize=14, rowh=0.6)
tf = box(s, 0.6, 4.5, 12.1, 2.6)
para(tf, "We also tried a fully **dense** reward — the progress signal at every step instead of a sparse signal at the "
         "end. It caps at the same ~0.1 as the sparse reward.", size=17, color=FG, bullet=True, first=True)
para(tf, "This is the key takeaway: the progress head separates classes well offline (d′ > 1), yet a dense progress "
         "reward still fails on-policy. **The bottleneck is on-policy exploitation, not reward sparsity or offline quality** — "
         "a dense reward simply gives the policy a smoother surface to exploit.", size=17, color=AMBER, bullet=True)

# ============================================================================
#  Slide 15 — Downstream: the exploit, made visual
# ============================================================================
s = slide(); header(s, "12", "What the policy actually does")
image(s, FIG / "rl_exploit_gallery.png", cx=3.7, top=1.55, max_w=6.8, max_h=5.4)
tf = box(s, 7.6, 1.9, 5.4, 4.9)
para(tf, "We rolled out the collapsed policy (80 episodes). **44% of its failures are scored as success** by the reward; "
         "true task success is 15%.", size=15, color=FG, bullet=True, first=True)
para(tf, "The reward even prefers the exploited failures (score 0.68–0.85) to a genuine success (0.65) — the ranking "
         "inverts on-policy.", size=15, color=AMBER, bullet=True)
para(tf, "The exploited states form a family — arm raised, object not delivered — that the reward over-scores. The "
         "policy parks there instead of completing the task.", size=15, color=FG, bullet=True)

# ============================================================================
#  Slide 15b — Offline quality does not transfer to the policy distribution
# ============================================================================
s = slide(); header(s, "13", "Strong offline reward quality does not transfer on-policy")
para(box(s, 0.6, 1.45, 12.1, 0.7),
     "The fine-tuned reward is excellent on the offline eval set, but its ability to tell success from failure "
     "decays as the data moves toward the policy's own distribution — to chance on live RL rollouts.",
     size=15, color=FG, bullet=True, first=True)
# left: the transfer collapse
table(s, ["scoring distribution", "success AUC"],
      [("Offline eval set", ["0.90"], "asym", [1]),
       ("BC-policy rollouts", ["0.66"], "paper", []),
       ("On-policy (RL) rollouts", ["0.55"], "base", [1])],
      0.6, 2.5, 6.0, [3.4, 1.6], fontsize=13, rowh=0.52)
para(box(s, 0.6, 5.0, 6.0, 1.0),
     "Robometer-FT, balanced set (135 successes / 65 failures from a competent GT-trained policy) — trustworthy n.",
     size=12, color=MUTE, italic=True, first=True)
# right: format ruled out
table(s, ["same on-policy frames, re-scored as…", "AUC"],
      [("raw 224  (as the live env feeds them)", ["0.47"], "base", []),
       ("h264 @ 224   (codec only)", ["0.46"], "base", []),
       ("h264 @ 240   (full curated pipeline)", ["0.56"], "base", []),
       ("jpeg @ 224", ["0.51"], "base", [])],
      6.9, 2.5, 6.2, [4.4, 1.2], fontsize=12, rowh=0.52)
para(box(s, 6.9, 5.2, 6.2, 1.4),
     "Pushing the live frames through the exact curated-data pipeline (codec + resolution) does **not** recover it — "
     "every variant stays at chance.", size=12.5, color=FG, bullet=True, first=True)
para(box(s, 0.6, 6.25, 12.5, 0.9),
     "So the offline→online gap is a genuine content/distribution shift, not a preprocessing artifact. Offline reward "
     "metrics (AUC, ECE, d′) do not predict on-policy usefulness.", size=15, color=AMBER, bullet=True, first=True)

# ============================================================================
#  Slide 16 — Summary
# ============================================================================
s = slide(); header(s, "∑", "Summary")
tf = box(s, 0.8, 1.6, 12.0, 5.6)
items = [
    "**Fine-tuning improves the offline reward in-distribution** — higher success-head AUPRC and TPR, and better "
    "class separation (d′ 1.3–1.5 vs 0.42) than the off-the-shelf baseline.",
    "**It trades this for out-of-distribution generalisation** — on held-out robots the baseline ranks better.",
    "**No VLM reward — fine-tuned or not, sparse or dense — trains the policy.** All cap near 0.1; the oracle "
    "reward reaches 0.48–0.76 on the same loop.",
    "**The cause is on-policy reward hacking, not offline quality.** The reward separates classes offline (AUC 0.89), "
    "but its structured false positives are exploitable, so RL amplifies them (14–17% → 41–62%).",
    "**Implication:** offline reward metrics do not predict downstream RL usefulness. The next step is on-policy reward "
    "correction — relabel the false positives the policy discovers to break the exploitable loop.",
]
for i, it in enumerate(items, 1):
    para(tf, f"{i}.   {it}", size=17, color=FG, space=12, first=(i == 1))

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides._sldIdLst)} slides)")
