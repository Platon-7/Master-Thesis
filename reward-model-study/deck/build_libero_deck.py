"""LIBERO progress-head diagnosis deck, styled after RL_progress_update_v3.pptx:
Arial, navy 1F3864 titles at 26 points with a thin accent bar, 13.5-point takeaway
boxes, 10-point gray footnotes, 13.33 x 7.5 inch slides."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from PIL import Image

DECK = "/shared/home/PKA4388/Master-Thesis/reward-model-study/deck"
FIGS = os.path.join(DECK, "figs_libero")
NAVY = RGBColor(0x1F, 0x38, 0x64)
GRAY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def txt(slide, x, y, w, h, text, size=13.5, bold=False, color=NAVY, align=PP_ALIGN.LEFT, line_spacing=1.12):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = line
        r.font.name = "Arial"; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return box

def header(slide, title, subtitle=None):
    t = title + ("\n" + subtitle if subtitle else "")
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.30), Inches(12.2), Inches(1.05))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.name = "Arial"; r.font.size = Pt(25); r.font.bold = True; r.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
        r2.font.name = "Arial"; r2.font.size = Pt(13); r2.font.bold = False; r2.font.color.rgb = GRAY
    bar = slide.shapes.add_shape(1, Inches(0.62), Inches(1.22), Inches(3.0), Inches(0.045))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

def pic(slide, path, y, max_w=12.0, max_h=4.6):
    im = Image.open(path); ar = im.width / im.height
    w = min(max_w, max_h * ar); h = w / ar
    x = (13.333 - w) / 2
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return y + h

def footnote(slide, text):
    txt(slide, 0.6, 7.02, 12.2, 0.45, text, size=10, color=GRAY)

# ---------------------------------------------------------------- slide 1: title
s = prs.slides.add_slide(BLANK)
txt(s, 0.9, 1.5, 11.5, 1.4, "The progress head on LIBERO:\nwhat breaks, why, and a fix that may avoid retraining",
    size=32, bold=True)
bar = s.shapes.add_shape(1, Inches(0.95), Inches(3.0), Inches(4.2), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
txt(s, 0.9, 3.4, 11.6, 3.0,
    "In one look:\n"
    "•  Our success head is now the strongest component: on live policy data it separates solved from failed episodes where the baseline's cannot.\n"
    "•  Our progress head is weaker than the baseline's on LIBERO: it reads even its own training demonstrations at 0.25 instead of 1.0.\n"
    "•  The cause is a two-peaked output distribution that the standard mean readout collapses. Reading the same output differently recovers the signal.\n"
    "•  Two issues in the authors' released code (a repeated-frame input and a changed training-target default) explain why the baseline looks stable while our models do not.",
    size=15, line_spacing=1.3)
txt(s, 0.9, 6.9, 11.5, 0.4, "Reward-model study, LIBERO task 28, July 3, 2026", size=11, color=GRAY)

# ------------------------------------------------- slide 2: heads on training clips
s = prs.slides.add_slide(BLANK)
header(s, "Scoring each model on the exact training clips of the evaluation task",
       "30 solved and 30 failed clips of 'close the top drawer of the cabinet', final-frame readings")
pic(s, f"{FIGS}/fig1_heads.png", 1.55, max_h=4.3)
txt(s, 0.9, 6.05, 11.6, 1.0,
    "Success head: failure data plus the asymmetric loss raised it from unusable (run3 reads true successes at 0.36) to strong (run2 at 0.82, failures at 0.20).\n"
    "Progress head: run2 reads demonstrations it trained on at 0.25. The defect is in the model output, not in the reinforcement learning pipeline.",
    size=13.5)
footnote(s, "Clips from the LIBERO-90 suite; they are inside our fine-tuning data and outside the baseline's, so this comparison favors us and still shows the 0.25 problem.")

# ------------------------------------------------------- slide 3: the bimodal shape
s = prs.slides.add_slide(BLANK)
header(s, "The progress head outputs a distribution; ours became two-peaked",
       "Average output on the solved training clips, final frame")
pic(s, f"{FIGS}/fig2_c51.png", 1.6, max_h=3.9)
txt(s, 0.9, 5.75, 11.6, 1.2,
    "The baseline puts its probability mass near the true value, so the mean is faithful (0.90).\n"
    "Both fine-tuned models park mass at zero while keeping a spike at 1.0. The mean of run2's distribution is 0.25, a value the model itself gives 3 percent probability.\n"
    "The information is present in the distribution; the mean readout throws it away.",
    size=13.5)
footnote(s, "All three models share the same 10-bin categorical progress head; only the training differs.")

# --------------------------------------------- slide 4: where the zero spike comes from
s = prs.slides.add_slide(BLANK)
header(s, "Where the zero spike comes from: the training targets, not the network")
pic(s, f"{FIGS}/fig3_windows.png", 1.5, max_h=3.1)
txt(s, 0.9, 4.85, 11.8, 2.1,
    "1.  A changed default in the authors' repository. The released baseline was trained with progress measured against the whole episode "
    "('absolute with respect to total frames'). The repository configuration we inherited measures progress inside each randomly cut training window "
    "('absolute first frame'): the same frame gets a different target depending on the cut, including exactly 0.0 whenever a window starts on it (figure above).\n"
    "2.  Our failure labels are heavy at zero: 52 percent of labeled failure frames map to target 0.0, and failures are oversampled.\n"
    "3.  The asymmetric loss amplifies the hedge: zero-bin mass grows from 0.42 (run3, standard loss) to 0.73 (run2, asymmetric loss).",
    size=13)
footnote(s, "Our per-frame failure labels are window-independent by construction and are not affected by issue 1; only the demonstration targets are.")

# ------------------------------------------------------------- slide 5: reward hacking
s = prs.slides.add_slide(BLANK)
header(s, "Why the baseline stays high and run2 decays: the margin at the top",
       "Reinforcement learning cares about the reward gap between 'almost solved' and 'solved'")
pic(s, f"{FIGS}/fig4_hacking.png", 1.55, max_h=4.0)
txt(s, 0.9, 5.75, 11.6, 1.1,
    "The baseline separates hovering from solving by 0.13 on a 0.9 scale; run2 by 0.05 on a 0.25 scale. Once the policy reaches the plateau, "
    "run2's gap disappears into critic noise and the entropy bonus: measured live, solved episodes collect 2.2 times the reward of failures early in training "
    "and only 1.1 times late. True success then decays while the reward stays high.",
    size=13.5)
footnote(s, "Left: measured on-policy values (mean peak progress of failed and solved live episodes). Right: one seed per curve, correct video input, no termination.")

# ---------------------------------------------------------- slide 6: repeated frame bug
s = prs.slides.add_slide(BLANK)
header(s, "A separate finding: the released code scores one repeated frame, not a video",
       "At every step the reward model receives the current frame repeated 8 or 16 times; real clips exist only in their real-robot pipeline")
pic(s, f"{FIGS}/fig5_framebug.png", 1.65, max_h=3.9)
txt(s, 0.9, 5.8, 11.6, 1.15,
    "The progress head survives because it effectively reads the state in the image. The success head must see the completion motion: on repeated frames it drops "
    "below the termination threshold on true successes, so success detection was silently disabled in every simulated run, ours and theirs.\n"
    "The bug even helps the baseline's training: a single-frame reward depends only on the current state, which is exactly what the learning algorithm assumes.",
    size=13.5)
footnote(s, "Confirmed in the unmodified upstream code and live logs (every scoring call received one distinct frame). The paper describes the opposite input format.")

# --------------------------------------------------------------- slide 7: top-bin fix
s = prs.slides.add_slide(BLANK)
header(s, "The proposed fix: read the top bin instead of the mean, no retraining",
       "The top bin is the model's stated probability that the task is complete")
pic(s, f"{FIGS}/fig6_readout.png", 1.6, max_h=3.9)
txt(s, 0.9, 5.7, 11.6, 1.15,
    "For run2 the top-bin readout separates solved from failed 3.3 to 1 (the mean gives 1.8 to 1). One line of inference code, same forward pass.\n"
    "The readout must match the model: the baseline's single-peaked head is best read by the mean (its top-bin run lost performance, 60 versus 90 percent peak), "
    "our two-peaked heads by the top bin. This symmetry is the honest framing, not 'our readout is better'.",
    size=13.5)
footnote(s, "Readouts computed from the same forward pass on the same clips as slide 3.")

# ------------------------------------------------- slide 8: success head on live data
s = prs.slides.add_slide(BLANK)
header(s, "The success head on live policy data: ours has an operating point, the baseline does not",
       "Source: full-horizon episodes from yesterday's runs (baseline: its no-termination run; run2: its max-reward run, in which the success vote never actually ended episodes). Per episode: peak success probability versus ground-truth end state.")
rows = [
    ["", "Separability (area under the curve)", "Calibrated threshold", "Detection rate", "False alarm rate"],
    ["run2 (ours)", "0.92", "0.85", "89 percent", "14 percent"],
    ["Robometer-4B (baseline)", "0.67", "0.965 (best possible)", "92 percent", "54 percent"],
]
tbl = s.shapes.add_table(3, 5, Inches(1.1), Inches(1.9), Inches(11.1), Inches(1.7)).table
tbl.columns[0].width = Inches(2.9)
for j in range(1, 5):
    tbl.columns[j].width = Inches(2.05)
for i, row in enumerate(rows):
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.name = "Arial"; r.font.size = Pt(13)
                r.font.bold = (i == 0) or (j == 0)
                r.font.color.rgb = NAVY if i else RGBColor(0xFF, 0xFF, 0xFF)
txt(s, 0.9, 4.1, 11.6, 1.9,
    "run2 separates solved from failed episodes well enough to operate: at threshold 0.85 it detects 89 percent of successes with 14 percent false alarms.\n"
    "For the baseline, solved and failed episodes both read about 0.97, so every threshold either misses successes or fires on half of the failures. "
    "This measured gap is the concrete value of training with failure data.\n\n"
    "We also found and fixed a second code issue: in the released simulation code the success vote never actually ended the episode "
    "(it only marked the stored transition as terminal, which corrupts the value targets). Our corrected version really terminates and resets.",
    size=13.5)
footnote(s, "run2: 9 solved and 215 failed episodes; baseline: 134 and 115. This table is the INPUT to the new true-termination runs launched today: run2 operates at 0.85 and the baseline at 0.965, each its own best threshold. Their curves are the next result, not shown here.")

# ----------------------------------------------------------- slide 9: early evidence
s = prs.slides.add_slide(BLANK)
header(s, "Early downstream evidence: both fixes were leading when the nodes were reclaimed",
       "Interrupted runs marked with a cross; all are repeating now on stable on-demand nodes")
pic(s, f"{FIGS}/fig7_early.png", 1.6, max_h=4.0)
txt(s, 0.9, 5.8, 11.6, 1.0,
    "run2 with the top-bin readout reached 100 percent at 50 thousand steps. run2 with success-head termination at the calibrated threshold reached "
    "100 percent already at 20 thousand steps, with 12 percent false terminations, matching the calibration's prediction.\n"
    "Rerunning now: run2 and run3 with the top-bin readout, and run2 and the baseline with success-head termination, each at its own calibrated threshold.",
    size=13.5)
footnote(s, "Single seed per configuration; the gray reference curve is the same run2 decay curve as slide 5.")

# ------------------------------------------------------------------ slide 10: summary
s = prs.slides.add_slide(BLANK)
header(s, "Where this leaves us")
txt(s, 0.9, 1.7, 11.8, 2.3,
    "Path A, no retraining (being tested right now)\n"
    "Top-bin readout for the progress reward, plus success-head termination at a calibrated threshold. "
    "If the reruns hold through 100 thousand steps, the LIBERO story is complete with inference-level fixes only.",
    size=14)
txt(s, 0.9, 3.6, 11.8, 2.2,
    "Path B, one more fine-tune (planned)\n"
    "Restore progress targets measured against the whole episode, keep the asymmetric loss only on the success head, "
    "and remove LIBERO-90 from the training data so it becomes a clean held-out evaluation with exactly the baseline's LIBERO exposure.",
    size=14)
txt(s, 0.9, 5.5, 11.8, 1.3,
    "Independent of both paths\n"
    "The repeated-frame input and the success vote that never terminates mean the simulated reinforcement learning results in the paper "
    "should be read with care; we now run corrected versions of both mechanisms.",
    size=14)
for y in (1.75, 3.65, 5.55):
    b = s.shapes.add_shape(1, Inches(0.62), Inches(y), Inches(0.06), Inches(0.55))
    b.fill.solid(); b.fill.fore_color.rgb = ORANGE if y == 3.65 else ACCENT; b.line.fill.background()

out = os.path.join(DECK, "LIBERO_progress_head_diagnosis.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
