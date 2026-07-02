"""Supervisor-facing results deck: progress since deck/VLM_reward_models.pptx.

4 slides, results-only, no internal jargon (no H1/H2, job ids, cFPR, debounce...).
Matches house style: 13.33 x 7.5 (16:9), Arial.
Sources: mega_results_288.csv / MEGA_VERDICT_288.md (288-run sweep) and
calib_matrix_{ft,4b}_full.json (15s/15f streaming calibration).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY  = RGBColor(0x1F, 0x38, 0x64)
BLUE  = RGBColor(0x2E, 0x5C, 0x9E)
GREEN = RGBColor(0x1E, 0x6B, 0x2E)
AMBER = RGBColor(0xA6, 0x60, 0x00)
RED   = RGBColor(0xB0, 0x1E, 0x2E)
GREY  = RGBColor(0x55, 0x55, 0x55)
HILITE= RGBColor(0xE6, 0xEE, 0xF6)
LGREY = RGBColor(0xF3, 0xF3, 0xF3)
CARD  = RGBColor(0xF4, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def setp(p, text, size=14, color=NAVY, bold=False, italic=False, space=6, align=None):
    p.text = text
    p.space_after = Pt(space)
    if align is not None:
        p.alignment = align
    if not p.runs:
        p.add_run()
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return p


def title(slide, text, sub=None):
    tf = tb(slide, 0.6, 0.38, 12.2, 1.1)
    setp(tf.paragraphs[0], text, size=28, bold=True, space=2)
    if sub:
        setp(tf.add_paragraph(), sub, size=13, color=GREY, italic=True, space=0)
    bar = slide.shapes.add_shape(1, Inches(0.62), Inches(1.22), Inches(3.0), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()


def table(slide, data, l, t, widths, row_h=0.36, fs=11.5,
          hi_rows=(), cell_colors=None, align_first_left=True):
    """data: list of row-tuples (row 0 = header). cell_colors: {(i,j): RGBColor}."""
    cell_colors = cell_colors or {}
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(l), Inches(t),
                                 Inches(sum(widths)), Inches(row_h * rows))
    tbl = shp.table
    tbl.first_row = False
    tbl.horz_banding = False
    for j, w in enumerate(widths):
        tbl.columns[j].width = Inches(w)
    for i in range(rows):
        for j in range(cols):
            c = tbl.cell(i, j)
            c.margin_top = c.margin_bottom = Pt(2)
            c.margin_left = Pt(5); c.margin_right = Pt(4)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0 and align_first_left) else PP_ALIGN.CENTER
            p.text = str(data[i][j])
            if not p.runs:
                p.add_run()
            r = p.runs[0]
            r.font.name = FONT
            r.font.size = Pt(fs)
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
                r.font.bold = True; r.font.color.rgb = WHITE
            else:
                c.fill.solid()
                c.fill.fore_color.rgb = HILITE if i in hi_rows else (WHITE if i % 2 else LGREY)
                r.font.color.rgb = NAVY
                if i in hi_rows:
                    r.font.bold = True
            if (i, j) in cell_colors:
                r.font.color.rgb = cell_colors[(i, j)]
                r.font.bold = True
    return tbl


# ============================ Slide 1 — headline ============================
s = prs.slides.add_slide(BLANK)
tf = tb(s, 0.8, 0.9, 11.7, 1.6)
setp(tf.paragraphs[0], "Autonomous RL with VLM Reward Models", size=36, bold=True, space=6)
setp(tf.add_paragraph(), "Results update — follows last week's reward-model deck", size=16, color=GREY, italic=True)
bar = s.shapes.add_shape(1, Inches(0.85), Inches(2.18), Inches(4.4), Pt(4))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

cards = [
    ("0.65 – 0.70", "true success rate of the best\nautonomous runs (CoffeePush)"),
    ("0", "ground-truth signals used\nduring RL training"),
    ("288", "runs in the controlled sweep\nthat isolates what matters"),
]
cx = 0.85
for big, cap in cards:
    card = s.shapes.add_shape(5, Inches(cx), Inches(2.85), Inches(3.7), Inches(2.1))
    card.fill.solid(); card.fill.fore_color.rgb = CARD; card.line.fill.background()
    ctf = card.text_frame; ctf.word_wrap = True
    ctf.margin_top = Pt(14); ctf.margin_left = Pt(14); ctf.margin_right = Pt(14)
    setp(ctf.paragraphs[0], big, size=38, bold=True, align=PP_ALIGN.CENTER, space=6)
    for line in cap.split("\n"):
        setp(ctf.add_paragraph(), line, size=12.5, color=GREY, align=PP_ALIGN.CENTER, space=0)
    cx += 3.95

tf = tb(s, 0.85, 5.45, 11.7, 1.6)
setp(tf.paragraphs[0],
     "Setting: the reward model watches the episode, provides the reward, and decides by itself when the task is "
     "solved — the episode ends on the model's decision, never on the simulator's success signal. Reported success "
     "is measured afterwards with ground truth, on separate evaluation episodes.",
     size=14, color=GREY, space=8)
setp(tf.add_paragraph(),
     "Last week: strong offline metrics, RL not yet working.   This week: autonomous RL works — under specific, "
     "now-understood conditions.",
     size=15, color=NAVY, bold=True)

# ============================ Slide 2 — the 288-run sweep ============================
s = prs.slides.add_slide(BLANK)
title(s, "Which design actually trains? — 288-run sweep on CoffeePush",
      "All 8 combinations of {reward model} × {which head detects success} × {which head provides reward}, "
      "everything else swept inside each row.")

data = [
    ("Reward model", "Detects success", "Provides reward", "Mean", "Best", "Runs ≥ 0.50"),
    ("Fine-tuned", "success head", "success head", "0.18", "0.70", "5 / 36"),
    ("Fine-tuned", "success head", "progress head", "0.07", "0.45", "0 / 36"),
    ("Fine-tuned", "progress head", "success head", "0.01", "0.05", "0 / 36"),
    ("Fine-tuned", "progress head", "progress head", "0.01", "0.40", "0 / 36"),
    ("Off-the-shelf", "success head", "success head", "0.01", "0.10", "0 / 36"),
    ("Off-the-shelf", "success head", "progress head", "0.04", "0.35", "0 / 36"),
    ("Off-the-shelf", "progress head", "success head", "0.01", "0.25", "0 / 36"),
    ("Off-the-shelf", "progress head", "progress head", "0.03", "0.25", "0 / 36"),
]
table(s, data, 0.6, 1.85, [1.55, 1.5, 1.5, 0.85, 0.85, 1.15],
      row_h=0.39, fs=11.5, hi_rows={1},
      cell_colors={(1, 3): GREEN, (1, 4): GREEN, (1, 5): GREEN})

tf = tb(s, 8.35, 1.85, 4.45, 4.9)
setp(tf.paragraphs[0], "What the winning row needs", size=15, bold=True, space=6)
for t in [
    "Fine-tuned reward model — the off-the-shelf model never reaches 0.50 anywhere (best 0.35).",
    "Success head for detection — ending episodes on the progress head fires mid-task and kills learning.",
    "Reward only at episode end (sparse). Per-step reward never produces a successful run.",
    "Detection threshold 0.80 – 0.85. At 0.90 the model never says “solved” → no reward → no learning.",
]:
    setp(tf.add_paragraph(), "•  " + t, size=12.5, color=GREY, space=7)
setp(tf.add_paragraph(),
     "Within the winning configuration, roughly 1 seed in 3 succeeds — the result is real but not yet stable.",
     size=12.5, color=NAVY, bold=True, space=0)

fn = tb(s, 0.6, 6.55, 12.2, 0.6)
setp(fn.paragraphs[0],
     "Each row aggregates 36 runs: 3 detection thresholds × reward at episode end vs every step × 1 vs 3 consecutive "
     "detections required × 3 seeds (8 × 36 = 288). “Mean / Best” = final true success; "
     "“Runs ≥ 0.50” = runs ending at or above 0.50.",
     size=10.5, color=GREY, italic=True, space=0)

# ============================ Slide 3 — calibration ============================
s = prs.slides.add_slide(BLANK)
title(s, "New: how reliable is the success detector on each task?",
      "Fine-tuned model, measured without RL: 15 successful + 15 failed scripted episodes per task; the detector "
      "watches each episode as it unfolds, exactly as it would in RL.")

cal = [
    ("Task", "Best threshold", "Detects true success", "False alarms on failures", "Fires before success", "Verdict"),
    ("CoffeePush", "0.85", "73 %", "27 %", "0 %", "Leaky but sufficient — RL trains here (proven)"),
    ("StickPull", "0.70", "87 %", "7 %", "20 %", "Good — best candidate for transfer"),
    ("BoxClose", "0.85", "100 %", "13 %", "60 % → 20 % with demo", "Workable if early firing is handled"),
    ("Assembly", "none", "—", "—", "—", "Unusable — success and failure overlap"),
]
table(s, cal, 0.6, 1.9, [1.5, 1.35, 1.85, 2.0, 2.1, 3.3], row_h=0.42, fs=11.5,
      cell_colors={(1, 5): GREEN, (2, 5): GREEN, (3, 5): AMBER, (4, 5): RED})

tf = tb(s, 0.6, 4.25, 12.2, 2.6)
setp(tf.paragraphs[0], "Two further findings", size=15, bold=True, space=6)
setp(tf.add_paragraph(),
     "•  The progress head cannot be used to decide termination on any task: it also rises on failed episodes "
     "(failure scores 0.3 – 0.7), so every threshold either fires constantly or never. As a reward magnitude it is "
     "fine (best run 0.45) — it must simply never end the episode.",
     size=12.5, color=GREY, space=7)
setp(tf.add_paragraph(),
     "•  Adding an in-context demonstration at inference is task-dependent: it fixes BoxClose’s early firing "
     "(60 % → 20 %) and improves StickPull, but degrades CoffeePush and Assembly — useful as a per-task option, "
     "not a global fix.",
     size=12.5, color=GREY, space=7)
setp(tf.add_paragraph(),
     "Rates are from 15 + 15 episodes per task (±13 %); read as ranges, not exact values.",
     size=10.5, color=GREY, italic=True, space=0)

# ============================ Slide 4 — in progress ============================
s = prs.slides.add_slide(BLANK)
title(s, "In progress", "All runs use the winning recipe; thresholds set per task and per model from the calibration above.")

items = [
    ("Transfer to harder tasks", "running",
     "54 runs: Assembly, BoxClose, StickPull × both reward models × 3 seeds. Calibration predicts StickPull "
     "transfers, BoxClose is borderline, Assembly fails.", "Results: TBD", BLUE),
    ("Fair re-test of the off-the-shelf model", "queued",
     "Its confidence scale sits far lower than the fine-tuned model’s (usable thresholds 0.15 – 0.70), so "
     "previous runs under-fired its detector. Re-running at its own calibrated thresholds — first on BoxClose, "
     "where its offline detection is strongest (AUC 0.94). On CoffeePush its detector is weak even at its own "
     "threshold, so the headline result stands either way.", "Results: TBD", BLUE),
    ("In-context demonstration during RL", "planned",
     "Calibration says a demonstration in context reduces early firing on BoxClose and helps StickPull — test it "
     "in the RL loop on those two tasks.", "Results: TBD", BLUE),
    ("Seed stability", "open question",
     "Only ~1 seed in 3 succeeds in the winning configuration. Understanding and closing this gap is the main "
     "open problem.", "", NAVY),
]
y = 1.7
for head, status, body, tbd, col in items:
    tf = tb(s, 0.7, y, 12.0, 1.25)
    p = setp(tf.paragraphs[0], f"{head}   ", size=15.5, bold=True, space=3)
    r = p.add_run(); r.text = f"[{status}]"
    r.font.name = FONT; r.font.size = Pt(12.5); r.font.italic = True; r.font.color.rgb = col
    setp(tf.add_paragraph(), body, size=12.5, color=GREY, space=2)
    if tbd:
        setp(tf.add_paragraph(), tbd, size=12, color=RED, bold=True, space=0)
    y += 1.38

out = "/shared/home/PKA4388/Master-Thesis/reward-model-study/RL_progress_update.pptx"
prs.save(out)
print("saved:", out, " slides:", len(prs.slides._sldIdLst))
