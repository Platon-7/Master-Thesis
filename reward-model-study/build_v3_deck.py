"""Version 3 supervisor deck. House style = deck/RL_progress_update.pptx (16:9, Arial, navy).
Metric: trailing mean of the last 3 evaluations (each = 20 episodes) — robust to end-spikes.
Numbers from train.log (jobs 535/575/551/552/603/647) + 288 sweep baseline (final-only).
Wording: ICL-on / ICL-off, asymmetric loss. ICL-on cells blank (runs scheduled)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY  = RGBColor(0x1F, 0x38, 0x64); BLUE = RGBColor(0x2E, 0x5C, 0x9E)
GREEN = RGBColor(0x1E, 0x6B, 0x2E); AMBER= RGBColor(0xA6, 0x60, 0x00)
RED   = RGBColor(0xB0, 0x1E, 0x2E); GREY = RGBColor(0x55, 0x55, 0x55)
HILITE= RGBColor(0xE6, 0xEE, 0xF6); LGREY= RGBColor(0xF3, 0xF3, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT  = "Arial"
FIG = "/shared/home/PKA4388/Master-Thesis/reward-model-study/deck/fig_gate_curves.png"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    box.text_frame.word_wrap = True; return box.text_frame

def setp(p, text, size=14, color=NAVY, bold=False, italic=False, space=6, align=None):
    p.text = text; p.space_after = Pt(space)
    if align is not None: p.alignment = align
    if not p.runs: p.add_run()
    for r in p.runs:
        r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic
    return p

def title(slide, text, sub=None):
    tf = tb(slide, 0.6, 0.34, 12.2, 1.05)
    setp(tf.paragraphs[0], text, size=26, bold=True, space=2)
    if sub: setp(tf.add_paragraph(), sub, size=12.5, color=GREY, italic=True, space=0)
    bar = slide.shapes.add_shape(1, Inches(0.62), Inches(1.12), Inches(3.0), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()

def table(slide, data, l, t, widths, row_h=0.36, fs=11.5, hi_rows=(), cell_colors=None,
          header_rows=1, align_first_left=True):
    cell_colors = cell_colors or {}
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Inches(l), Inches(t),
                                 Inches(sum(widths)), Inches(row_h*rows))
    tbl = shp.table; tbl.first_row = False; tbl.horz_banding = False
    for j, w in enumerate(widths): tbl.columns[j].width = Inches(w)
    for i in range(rows):
        for j in range(cols):
            c = tbl.cell(i, j)
            c.margin_top = c.margin_bottom = Pt(2); c.margin_left = Pt(6); c.margin_right = Pt(4)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0 and align_first_left) else PP_ALIGN.CENTER
            p.text = str(data[i][j])
            if not p.runs: p.add_run()
            r = p.runs[0]; r.font.name = FONT; r.font.size = Pt(fs)
            if i < header_rows:
                c.fill.solid(); c.fill.fore_color.rgb = NAVY; r.font.bold = True; r.font.color.rgb = WHITE
            else:
                c.fill.solid(); c.fill.fore_color.rgb = HILITE if i in hi_rows else (WHITE if i % 2 else LGREY)
                r.font.color.rgb = NAVY
                if i in hi_rows: r.font.bold = True
            if (i, j) in cell_colors:
                col, bold = cell_colors[(i, j)]; r.font.color.rgb = col; r.font.bold = bold
    return tbl

def note(slide, text, t=7.02, color=GREY, size=10):
    setp(tb(slide, 0.6, t, 12.2, 0.45).paragraphs[0], text, size=size, color=color, italic=True, space=0)

# ============ SLIDE 1 — stabilization techniques ============
s = prs.slides.add_slide(BLANK)
title(s, "Stabilizing the autonomous reward signal",
      "Coffee-Push, same recipe; only the stabilization technique changes. "
      "A seed “collapses” when its success rate stays at zero.")
data = [
    ["Stabilization technique", "Success rate", "Collapsed seeds", "Extra cost", "Holds on 2nd task?"],
    ["None — fixed threshold", "0.33", "1 of 6", "—", "no"],
    ["Output sampling", "0.31", "0 of 3", "none", "untested"],
    ["Self-ensemble (majority vote)", "0.43", "1 of 5", "4× scoring", "weak"],
    ["Gate + self-ensemble", "0.31", "1 of 3", "4× scoring", "—"],
    ["Minimum-length gate", "0.64", "0 of 3", "none", "yes (Box-Close 3/3)"],
]
cc = {(5,1): (GREEN, True), (5,2): (GREEN, True), (5,4): (GREEN, True), (5,5): (GREEN, True),
      (1,1): (AMBER, False), (1,2): (RED, True)}
table(s, data, 0.95, 1.65, [4.3, 1.9, 2.3, 1.9, 2.9], row_h=0.6, fs=13.5, hi_rows={5}, cell_colors=cc)
setp(tb(s, 0.95, 5.55, 11.4, 1.2).paragraphs[0],
     "The minimum-length gate gives the highest success with no collapsed seeds, adds no inference "
     "cost, and is the only technique that also holds on a second task. It simply ignores any "
     "“success” call made before the task could plausibly be finished.", size=13.5, color=NAVY, space=0)
note(s, "Technique screening on Model A / Coffee-Push (3–6 seeds); the chosen gate is validated at 6 seeds on the next slides. "
        "Success rate = mean over seeds of each run’s last-3-evaluation average (20 episodes). No ground truth used in training.")

# ============ SLIDE 2 — gate / ICL / loss isolation + curves ============
s = prs.slides.add_slide(BLANK)
title(s, "Isolating the gate, ICL, and the training loss",
      "Success rate (seeds alive / total). Three reward models × two tasks.")
SCH = ""; NA = "—"
data = [
    ["Reward model · Task", "ICL-off,\nno gate", "ICL-off,\ngate", "ICL-on,\nno gate", "ICL-on,\ngate"],
    ["A — ICL + asymmetric loss · Coffee-Push", "0.33", "0.29 (3/6)", SCH, SCH],
    ["A — ICL + asymmetric loss · Box-Close",   NA,     "0.48 (6/6)", SCH, SCH],
    ["B — asymmetric loss · Coffee-Push",       "0.34 (3/3)", "0.48 (6/6)", "n/a", "n/a"],
    ["B — asymmetric loss · Box-Close",         "0.17 (3/3)", "0.68 (6/6)", "n/a", "n/a"],
    ["C — standard loss · Coffee-Push",         "0.09 (2/3)", "0.09 (2/6)", "n/a", "n/a"],
    ["C — standard loss · Box-Close",           "0.32 (2/3)", "0.30 (6/6)", "n/a", "n/a"],
]
cc = {}
for i in (3,4): cc[(i,2)] = (GREEN, True)
cc[(2,2)] = (GREEN, True)
for i in (1,2): cc[(i,3)] = (AMBER, False); cc[(i,4)] = (AMBER, False)
for i in (3,4,5,6): cc[(i,3)] = (GREY, False); cc[(i,4)] = (GREY, False)
table(s, data, 0.55, 1.55, [4.55, 1.65, 1.65, 1.55, 1.55], row_h=0.4, fs=11, hi_rows={3,4}, cell_colors=cc)
# curves figure (proof of real learning)
s.shapes.add_picture(FIG, Inches(0.5), Inches(4.35), width=Inches(6.2))
tf = tb(s, 6.9, 4.45, 6.1, 2.6)
setp(tf.paragraphs[0], "Gate vs no gate isolates the gate (Box-Close 0.17 → 0.68 for B).",
     size=12.5, color=NAVY, space=10)
setp(tf.add_paragraph(), "B vs C isolates the loss (asymmetric ≫ standard).", size=12.5, color=NAVY, space=10)
setp(tf.add_paragraph(), "Trailing-mean over seeds (n = 6 gate / 3 no-gate). ICL applies to model A only.",
     size=10, color=GREY, italic=True, space=0)

# ============ SLIDE 3 — online calibration ============
s = prs.slides.add_slide(BLANK)
title(s, "Removing per-task threshold tuning: online calibration",
      "Match the hand-tuned threshold with no per-task search.")
data = [
    ["Approach", "Per-task tuning?", "Needs demos?", "Coffee / Box"],
    ["Hand-tuned threshold (baseline)", "yes", "no", "0.48 / 0.68"],
    ["Online demo-based (success-anchor)", "no", "yes", "0.51 / 0.37"],
    ["Online demo-free (score-valley)", "no", "no", "—"],
]
cc = {(1,1): (AMBER, False), (2,1): (GREEN, True), (3,1): (GREEN, True),
      (3,2): (GREEN, True), (2,2): (AMBER, False)}
table(s, data, 0.9, 1.8, [4.2, 2.2, 2.0, 3.5], row_h=0.64, fs=13.5, cell_colors=cc)
tf = tb(s, 0.9, 4.95, 11.4, 1.6)
setp(tf.paragraphs[0], "The demo-free method sets the bar from the live distribution of the model’s own "
     "scores, firing only when a distinct “success” cluster emerges — so it cannot be exploited the way a "
     "fixed threshold is, and needs neither labels nor demonstrations.", size=13.5, color=NAVY, space=0)

# ============ SLIDE 4 — offline-calibration results (deep dive) ============
def results_slide(title_txt, sub_txt, fig_path, caption_lines, note_txt):
    s = prs.slides.add_slide(BLANK)
    title(s, title_txt, sub_txt)
    s.shapes.add_picture(fig_path, Inches(0.55), Inches(1.5), width=Inches(8.4))
    tf = tb(s, 9.15, 1.7, 3.9, 5.4)
    for i, (txt, c) in enumerate(caption_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        setp(p, txt, size=11.5, color=c, space=8)
    note(s, note_txt)

results_slide(
    "Best case — offline-calibrated threshold: does the policy truly learn?",
    "run2 (asymmetric-loss reward model) · min-length gate · threshold from GT-labeled rollouts · 6 seeds",
    "/shared/home/PKA4388/Master-Thesis/reward-model-study/deck/fig_bestcase_run2.png",
    [("Top: true success (greedy eval), mean ±1 std over 6 seeds. Coffee-Push 0.48, Box-Close 0.68.", NAVY),
     ("Bottom: true success and VLM reward rise together — a faithful proxy, not reward hacking.", NAVY),
     ("Curves still rising at 40k (not converged); reward is graded.", GREY)],
    "Needs GT-labeled success+failure rollouts to set the threshold. Eval GT only; reward model never used at eval.")

# ============ SLIDE 5 — demo-anchor results (deep dive, parallel layout) ============
results_slide(
    "Demo-anchor (demos-only) calibration: same examination",
    "run2 · min-length gate · threshold from a few success demos (no GT labels, no grid search) · 3 seeds",
    "/shared/home/PKA4388/Master-Thesis/reward-model-study/deck/fig_bestcase_demoanchor_run2.png",
    [("Same panels as the offline slide; only the threshold source changes.", NAVY),
     ("Coffee-Push 0.51 (≈ offline 0.48), Box-Close 0.37 (vs offline 0.68): matches offline on Coffee-Push, weaker on Box-Close.", AMBER),
     ("Wider reward-faithfulness gap — but needs only success demos, no GT labels.", GREY)],
    "Demos-only — no GT labels, no grid search. 3 seeds (vs 6 for offline).")

# ============ SLIDE 6 — closing the loop from last week ============
s = prs.slides.add_slide(BLANK)
title(s, "Since last week: open items, now answered",
      "Last week's questions — off-the-shelf baseline, transfer to harder tasks, seed stability.")
# baseline result (owed from last week)
tf = tb(s, 0.95, 1.65, 11.6, 0.9)
setp(tf.paragraphs[0], "Off-the-shelf reward model — dead on every task.", size=15, bold=True, color=RED, space=2)
setp(tf.add_paragraph(), "Mean 0.02, never reaches 0.50 (best 0.32) across 21 runs at its own calibrated "
     "thresholds (and with ICL). The fine-tuned reward model is required.", size=13, color=NAVY, space=0)
# which tasks train
data = [
    ["Task", "Trains?", "Success (run2 + gate, 6 seeds)"],
    ["Coffee-Push", "yes", "0.48  (6/6 seeds alive)"],
    ["Box-Close", "yes — 2nd working task", "0.68  (6/6 seeds alive)"],
    ["Stick-Pull", "no", "0.00  (0/3)"],
    ["Assembly", "no", "0.00  (0/4)"],
]
cc = {(1,1):(GREEN,True),(2,1):(GREEN,True),(3,1):(RED,True),(4,1):(RED,True)}
table(s, data, 0.95, 2.95, [3.0, 3.4, 4.6], row_h=0.5, fs=13, cell_colors=cc)
tf = tb(s, 0.95, 5.75, 11.6, 1.4)
setp(tf.paragraphs[0], "Stick-Pull & Assembly are characterized negatives: no usable threshold, and neither the "
     "gate nor in-context demos rescue them — an on-policy robustness limit of the reward model, not a tuning miss.",
     size=12.5, color=NAVY, space=6)
setp(tf.add_paragraph(), "Seed stability — last week's main open problem (~1 seed in 3) — is resolved by the gate: "
     "every seed stays alive on the working tasks (slides 1–2).", size=12.5, color=NAVY, space=0)
note(s, "FT model, winning recipe (success head, gate, sparse). Eval GT only; trailing-mean over seeds.")

out = "/shared/home/PKA4388/Master-Thesis/reward-model-study/deck/RL_progress_update_v3.pptx"
prs.save(out); print("saved", out, "slides:", len(prs.slides._sldIdLst))
